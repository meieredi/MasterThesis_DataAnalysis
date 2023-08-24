import os
import numpy as np
import matplotlib as mpl
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

def plotDLS(corrData, intData, fileNames, param):
    """
    Plot DLS correlation (and intensity) curve(s) in a 1x2 (w/ intensity data) 
    or 1x1 (w/o intensity data) plot matrix
    """
    
    # Extract plotting parameters
    [autoPlot, plotIntData, avgPlot, plotNumberOfExp, plotFileName, plotTitle, plotLegend, colors, cmap, font, showTitle] = param

    # Preliminary parameters
    corrLimX = [0.5, 1e5]
    corrLimY = [-0.05, 1.05]
    intLimX = [0.4e-3, 10]
    intLimY = [-2,27]
    
    corrTextCoord = [2*corrLimX[0], 0.92*corrLimY[1]]
    intTextCoord = [2*intLimX[0], 0.92*intLimY[1]] 
    
    text = 'Number of Experiments: {:.0f}'
    
    multipleFiles = False
    
    # First time parameter (true only for first loop iteration)
    firstTime = True
        
    # Array to save whether intensity data is stored in each given file
    intMeasured = np.zeros(len(fileNames))
    
    # Preliminary check for intensity data
    for file, i in zip(fileNames, range(len(fileNames))):

        if intData[file].any():
            # Set bool to true
            intMeasured[i] = 1
         
    if len(fileNames) > 1:
        multipleFiles = True

    
    # Loop over all given files
    for i, file in enumerate(fileNames):
        
        # Check shape of correlation data array
        corrShape = np.shape(corrData[file])
        
        # Calculate number of experiments (one triplicate per experiment)
        corrExpNum = int((corrShape[1] - 1)/3)
          
        # Adjust colormap settings
        if multipleFiles:
            if len(fileNames)%2 == 1:
                darkGray = [0.45, 0.45, 0.45, 1]  # RGBA value, in closed interval [0,1]
                colMap = getattr(plt.cm, cmap)(np.linspace(0,1,len(fileNames)))
                colMap[1] = darkGray
                plt.rcParams["axes.prop_cycle"] = plt.cycler("color", colMap)
            else:
                plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,len(fileNames))))
        else:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,corrExpNum)))

        # Extract correlation data triplets (and compute row-wise mean & std. dev.)
        lagTime = corrData[file][:, 0]
        
        meanExpCorr = np.empty((corrShape[0], int(corrExpNum)))
        stdExpCorr = np.empty((corrShape[0], int(corrExpNum)))
        
        for j in range(int(corrExpNum)): 
            meanExpCorr[:,j] = np.mean(corrData[file][:, (3*j+1):(3*j+4)], axis=1)
            stdExpCorr[:,j] = np.std(corrData[file][:, (3*j+1):(3*j+4)], axis=1, ddof=1)
            
        if avgPlot and corrExpNum > 1:
            meanCorr = np.mean(meanExpCorr, axis=1)
            stdCorr = np.std(meanExpCorr, axis=1)
            
        elif avgPlot and corrExpNum == 1:
            meanCorr = meanExpCorr[:,0]
            stdCorr = stdExpCorr[:,0]
        
        else:
            meanCorr = meanExpCorr
            stdCorr = stdExpCorr

        # Check whether current file contains intensity data
        if intMeasured.any() and plotIntData:
         
            # Check  shape of intensity data array
            intShape = np.shape(intData[file])
            
            # Calculate number of experiments (one triplicate per experiment)
            if intShape[1] > 1:
                intExpNum = (intShape[1] - 1)/3
            else: 
                intExpNum = 1.0
                
            if firstTime:
                # Initialize 1x2 subplot matrix
                fig, axs = plt.subplots(nrows=1, ncols=2, figsize=(14, 5))
                
                # Set axis labels and legends
                axs[0].set_xlabel('Lag Time [$\mu$s]')
                axs[0].set_ylabel('Correlation Coefficient [-]')
                axs[0].set_xlim(corrLimX)
                axs[0].set_ylim(corrLimY)
                axs[0].tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                     bottom=True, top=False, left=True, right=True, direction='in')
                axs[1].set_xlabel('Hydrodynamic Diameter [$\mu$m]')
                axs[1].set_ylabel('Intensity [\%]') if font == 'LaTeX' else axs[1].set_ylabel('Intensity [%]')                    
                axs[1].set_xlim(intLimX)
                axs[1].set_ylim(intLimY)
                axs[1].tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                     bottom=True, top=False, left=True, right=True, direction='in')
                
                if multipleFiles and plotNumberOfExp:
                    axs[0].text(corrTextCoord[0], corrTextCoord[1], text.format(len(fileNames)), va='center')
                    axs[1].text(intTextCoord[0], intTextCoord[1], text.format(np.count_nonzero(intMeasured)), va='center')
                elif plotNumberOfExp:
                    axs[0].text(corrTextCoord[0], corrTextCoord[1], text.format(corrExpNum), va='center')
                    axs[1].text(intTextCoord[0], intTextCoord[1], text.format(intExpNum), va='center')

                firstTime = False
                
            if avgPlot:
                # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
                if colors[0]:
                    axs[0].semilogx(lagTime, meanCorr, color=colors[i])
                    axs[0].fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, color=colors[i], alpha=0.25,
                                        linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)
                else:
                    axs[0].semilogx(lagTime, meanCorr)
                    axs[0].fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, alpha=0.25,
                        linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)
                    
                axs[0].xaxis.get_major_locator().set_params(numticks=99)
                axs[0].xaxis.get_minor_locator().set_params(numticks=99, subs='auto')
                
            
            else:
                for j in range(int(corrExpNum)): 
                    # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
                    if colors[0]:
                        axs[0].semilogx(lagTime, meanCorr[:,j], color=colors[0])
                        axs[0].fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                            color=colors[0], alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                            antialiased=True)
                    else:
                        if multipleFiles:
                            fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                            axs[0].semilogx(lagTime, meanCorr[:,j], color=fileColor)
                            axs[0].fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                                color=fileColor, alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                                antialiased=True)
                        else:
                            axs[0].semilogx(lagTime, meanCorr[:,j])
                            axs[0].fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                                alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                                antialiased=True)
                        
                    axs[0].xaxis.get_major_locator().set_params(numticks=99)
                    axs[0].xaxis.get_minor_locator().set_params(numticks=99, subs='auto')

            
            # Access color of previous curve
            prevCol = axs[0].get_lines()[-1].get_color()
            
            if intMeasured[i]:
                # Extract intensity data triplets (and compute mean & std. dev.)
                size = intData[file][:, 0]
                
                if avgPlot:
                    meanInt = np.mean(intData[file][:, 1:], axis=1)
                    stdInt = np.std(intData[file][:, 1:], axis=1, ddof=1)
                
                else:
                    meanInt = np.empty((intShape[0], int(intExpNum)))
                    stdInt = np.empty((intShape[0], int(intExpNum)))
                    
                    for j in range(int(intExpNum)): 
                        meanInt[:,j] = np.mean(intData[file][:, (3*j+1):(3*j+4)], axis=1)
                        stdInt[:,j] = np.std(intData[file][:, (3*j+1):(3*j+4)], axis=1, ddof=1)
            
                if avgPlot:
                    # Add curves to subfigure 2 (plot intensity [%] vs. size [um])
                    axs[1].semilogx(size/1000, meanInt, color=prevCol, subs='auto')
                    axs[1].fill_between(size/1000, meanInt-stdInt, meanInt+stdInt, color=prevCol, alpha=0.25,
                    linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)
                    axs[1].xaxis.get_major_locator().set_params(numticks=99)
                    axs[1].xaxis.get_minor_locator().set_params(numticks=99, subs='auto')


                    
                else:
                    for j in range(int(intExpNum)): 
                        # Add curves to subfigure 2 (plot intensity [%] vs. size [um])
                        if colors[0]:
                            axs[1].semilogx(size/1000, meanInt[:,j], subs='auto', color=colors[0])
                            axs[1].fill_between(size/1000, meanInt[:,j]-stdInt[:,j], meanInt[:,j]+stdInt[:,j], 
                            color=colors[0], alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                            antialiased=True)
                        else:
                            if multipleFiles:
                                fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                                axs[1].semilogx(size/1000, meanInt[:,j], subs='auto', color=fileColor)
                                axs[1].fill_between(size/1000, meanInt[:,j]-stdInt[:,j], meanInt[:,j]+stdInt[:,j], 
                                color=fileColor, alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                                antialiased=True)
                            else:
                                axs[1].semilogx(size/1000, meanInt[:,j], subs='auto')
                                axs[1].fill_between(size/1000, meanInt[:,j]-stdInt[:,j], meanInt[:,j]+stdInt[:,j], 
                                alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                                antialiased=True)
                            
                        axs[1].xaxis.get_major_locator().set_params(numticks=99)
                        axs[1].xaxis.get_minor_locator().set_params(numticks=99, subs='auto')

            if plotLegend:

                if plotFileName == 'JJn_01mg_ProtType_PM' or plotFileName == 'JJn_01mg_ProtType_ST' \
                    or plotFileName == 'JJ1_01mg_BGr_PM' or plotFileName == 'JJ1_01mg_BGr_ST':
                    axs[0].legend(plotLegend, fontsize='small')
                else:
                    axs[0].legend(plotLegend)
                #axs[1].legend()

                if not avgPlot:
                    for i in range(len(fileNames)):
                        leg = axs[0].get_legend()
                        fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                        leg.legendHandles[i].set_color(fileColor)
            
        # No intensity data found
        else:
        
            if firstTime:
                # Create 1x1 subplot matrix
                if autoPlot:
                    fig, axs = plt.subplots(nrows=1, ncols=1, figsize=(5.5, 5))
                else:
                    fig, axs = plt.subplots(nrows=1, ncols=1, figsize=(7, 5))

                # Set axis labels and legends
                axs.set_xlabel('Lag Time [$\mu$s]')
                axs.set_ylabel('Correlation Coefficient [-]')
                axs.set_xlim(corrLimX)
                axs.set_ylim(corrLimY)
                axs.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                                    bottom=True, top=False, left=True, right=True, direction='in')
                if multipleFiles and plotNumberOfExp:
                    axs.text(corrTextCoord[0], corrTextCoord[1], text.format(len(fileNames)), va='center')
                elif plotNumberOfExp:
                    axs.text(corrTextCoord[0], corrTextCoord[1], text.format(corrExpNum), va='center')
                
                firstTime = False
            
            if avgPlot:
                # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
                if colors[0]:
                    axs.semilogx(lagTime, meanCorr, color=colors[i])
                    axs.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, color=colors[i], alpha=0.25,
                    linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)
                else:
                    axs.semilogx(lagTime, meanCorr)
                    axs.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, alpha=0.25,
                    linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)
                    
                axs.xaxis.get_major_locator().set_params(numticks=99)
                axs.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')

            
            else:
                for j in range(int(corrExpNum)): 
                    # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
                    if colors[0]:
                        axs.semilogx(lagTime, meanCorr[:,j], color=colors[i])
                        axs.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        color=colors[i], alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)
                    else:
                        if multipleFiles:
                            fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                            axs.semilogx(lagTime, meanCorr[:,j], color=fileColor)
                            axs.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                            color=fileColor, alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                            antialiased=True)
                        else:
                            axs.semilogx(lagTime, meanCorr[:,j])
                            axs.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                            alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                            antialiased=True)

                    axs.xaxis.get_major_locator().set_params(numticks=99)
                    axs.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')


            if plotLegend and i == (len(fileNames) - 1):
                
                if plotFileName == 'JJn_01mg_ProtType_PM' or plotFileName == 'JJn_01mg_ProtType_ST' \
                    or plotFileName == 'JJ1_01mg_BGr_PM' or plotFileName == 'JJ1_01mg_BGr_ST':
                    axs.legend(plotLegend, fontsize='small')
                else:
                    axs.legend(plotLegend)

                if not avgPlot:
                    for i in range(len(fileNames)):
                        leg = axs.get_legend()
                        fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                        leg.legendHandles[i].set_color(fileColor)
    
    # Automatically detect experimental conditions (if autoPlot is enabled)
    # REMARK: Currently only detects conditions of first file in given fileList!
    
    if autoPlot:
        plotTitle, exportNamePdf, exportNamePng = scanExpCond(fileNames, 0, 0, avgPlot, font, IsItDLS = 1)
    
    # Set export title if autoPlot is not selected
    else:
        endNamePdf = plotFileName + '.pdf'
        endNamePng = plotFileName + '.png'

        exportNamePdf = os.path.join('01_DLS', '01_Overview','01_PDF', endNamePdf)
        exportNamePng = os.path.join('01_DLS', '01_Overview','02_PNG', endNamePng)
                                     
    # Set overall title
    if showTitle:
        if intMeasured.any():
            fig.suptitle(plotTitle, y=0.97)
        else:
            fig.suptitle(plotTitle, y=0.97)
        
    # Save plot    
    plt.savefig(os.path.join('02_Plots', exportNamePdf), bbox_inches='tight')
    plt.savefig(os.path.join('02_Plots', exportNamePng), bbox_inches='tight')
    
    # Don't show plot in jupyter notebook
    plt.close()

    return


def plotMDI(countData, fileNames, param):
    """
    Plot MDI counts in a grouped bar plot
    """
    [autoPlot, rmBin, logPlot, avgPlot, plotLimits, expNames, plotFileName, plotTitle, colors, cmap, font, showTitle] = param

    # Preliminary parameters
    if plotLimits:
        countLimY = plotLimits
        
    # Tuple containing experiment number strings
    if expNames[0]:
        expTuple = expNames
    else:
        expTuple = ()
    
    # Set multipleFiles parameter to zero by default
    multipleFiles = False
    numFiles = len(fileNames)
        
    # First time parameter (true only for first loop iteration)
    firstTime = True
    
    # Initialize empty count dictionary (needed for grouped barplot)
    if rmBin:
        keyList = ['Total ($\geq$ 5 $\mu$m)', '5-10 $\mu$m', '10-25 $\mu$m', '25-70 $\mu$m', '> 70 $\mu$m']
    else:
        keyList = ['Total ($\geq$ 1 $\mu$m)', '1-2 $\mu$m', '2-5 $\mu$m', '5-10 $\mu$m', '10-25 $\mu$m']
    keyLength = len(keyList)
    keyContent = []
    meanCountDict = dict.fromkeys(keyList, keyContent)
    stdCountDict = dict.fromkeys(keyList, keyContent)
        
    # Check whether multiple filenames are given
    if numFiles > 1:
        multipleFiles = True

    # Adjust colormap settings
    if rmBin:
        lightGray = [0.7, 0.7, 0.7, 1]  # RGBA value, in closed interval [0,1]
        colMap = getattr(plt.cm, cmap)(np.linspace(0,1,6))#len(keyList)))
        plt.rcParams["axes.prop_cycle"] = plt.cycler("color", np.insert(colMap, 2, lightGray, axis=0))
    else:
        darkGray = [0.4, 0.4, 0.4, 1]   # RGBA value, in closed interval [0,1]
        colMap = getattr(plt.cm, cmap)(np.linspace(0,1,6))#len(keyList)))
        plt.rcParams["axes.prop_cycle"] = plt.cycler("color", np.insert(colMap, 0, darkGray, axis=0))

    countDataProc = {}
    countDataDel = {}

    # Loop over all given files
    for i, file in enumerate(fileNames):
        
        
        # Remove first bin
        if rmBin:
            countDataDel[file] = np.delete(countData[file], [0,1], 0)
            countDataProc[file] = np.insert(countDataDel[file], 0, np.sum(countDataDel[file], axis=0), axis=0)

        else:
            countDataProc[file] = np.insert(countData[file], 0, np.sum(countData[file], axis=0), axis=0)
            #countDataProc[file] = countData[file]

        
        # Find max. count in given file
        countMax = np.max(countDataProc[file])
        
        # Check shape of correlation data array
        countShape = np.shape(countDataProc[file])
        
        # Calculate number of experiments (one duplicate per experiment)
        countExpNum = int(countShape[1] / 2)
            
        # Attention: first exp. file must contain more measurements than any later file (of same exp. cond.), otherwise: error!
        # In general: maximum 3 files per experimental condition!
        if firstTime:
            meanExpCount = np.zeros((countShape[0], countExpNum*numFiles))
            stdExpCount = np.zeros((countShape[0], countExpNum*numFiles)) 
            
        for j in range(int(countExpNum)):
            meanExpCount[:,i*countExpNum+j] = np.mean(countDataProc[file][:, (2*j):(2*j+2)], axis=1)
            stdExpCount[:,i*countExpNum+j] = np.std(countDataProc[file][:, (2*j):(2*j+2)], axis=1)#, ddof=1)
            
        if avgPlot:
            
            if firstTime:
                meanCount = np.zeros((countShape[0], numFiles))
                stdCount = np.zeros((countShape[0], numFiles))
            
            if i == 0:
                meanCount[:,i] = np.mean(meanExpCount[:,:(i*countExpNum+countExpNum)], axis=1)
                stdCount[:,i] = np.std(meanExpCount[:,:(i*countExpNum+countExpNum)], axis=1)#, ddof=1)
            else: 
                meanCount[:,i] = np.mean(meanExpCount[:,((i-1)*countExpNum+countExpNum):(i*countExpNum+countExpNum)], axis=1)
                stdCount[:,i] = np.std(meanExpCount[:,((i-1)*countExpNum+countExpNum):(i*countExpNum+countExpNum)], axis=1)#, ddof=1)

            if not expNames[0]:
                addTuple = 'Mean ' + str(i+1)
                expTuple += (addTuple,)
        
        else: 
            
            meanCount = meanExpCount
            stdCount = stdExpCount
              
            for j in range(int(countExpNum)):

                if not expNames[0]:
                    addTuple = 'Exp.' + str(i+1) + '.' + str(j+1)
                    expTuple += (addTuple,)
                
        
        if firstTime:
            
            if avgPlot:
                fig, ax = plt.subplots(layout='constrained', figsize=(4+3*numFiles/3, 5))
                
            else:
                fig, ax = plt.subplots(layout='constrained', figsize=(4+3*countExpNum*numFiles/3, 5))

            if rmBin:
                next(ax._get_lines.prop_cycler)
                next(ax._get_lines.prop_cycler)

            # Add some text for labels, title and custom x-axis tick labels, etc.
            #ax.set_xlabel('Experiments')
            ax.set_ylabel('Number Density [mL$^{-1}$]')
            ax.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                            bottom=True, top=False, left=True, right=True, direction='in')

            firstTime = False

            if rmBin:
                ax.bar([], [], label='_nolegend_')
                ax.bar([], [], label='_nolegend_')

    x = np.arange(len(expTuple))
    width = 0.1  # the width of the bars
    multiplier = 1
    
    if avgPlot:
        
        for j in range(keyLength):                
            meanCountDict[keyList[j]] = meanCount[j,:]
            stdCountDict[keyList[j]] = stdCount[j,:]
          
        for attribute, measurement in meanCountDict.items():
            offset = width * multiplier
            rects = ax.bar(x + offset, measurement, width, label=attribute, alpha=0.7)
            ax.errorbar(x + offset, measurement, yerr=stdCountDict[attribute], ecolor='black',
                       capsize=5, linestyle='')
            #ax.bar_label(rects, label_type='center', padding=3)
            multiplier += 1
            
    else:
        
        for j in range(keyLength):                
            meanCountDict[keyList[j]] = meanCount[j,:]
            stdCountDict[keyList[j]] = stdCount[j,:]
          
        for attribute, measurement in meanCountDict.items():
            offset = width * multiplier
            rects = ax.bar(x + offset, measurement, width, label=attribute, alpha=0.7)
            ax.errorbar(x + offset, measurement, yerr=stdCountDict[attribute], ecolor='black',
                       capsize=5, linestyle='')
            #ax.bar_label(rects, label_type='center', padding=3)
            multiplier += 1            

    # Set additional plotting parameters
    ax.set_xticks(x + (keyLength+1)/2*width, expTuple)
    ax.legend(ncols=1)
    
    if logPlot:
        ax.set_yscale('log')
    else:
        ax.set_yscale('linear')

    if plotLimits and logPlot:
        ax.set_ylim(countLimY)
    elif plotLimits:
        ax.set_ylim(countLimY)
        ax.ticklabel_format(style='sci', axis='y', scilimits=(0,0), useMathText=True)
    elif logPlot:
        ax.set_ylim(bottom=10)
    else:
        ax.set_ylim(bottom=0)
        ax.ticklabel_format(style='sci', axis='y', scilimits=(0,0), useMathText=True)


    # Automatically detect experimental conditions (if autoPlot is enabled)
    # REMARK: Currently only detects conditions of first file in given fileList!
    
    if autoPlot:
        plotTitle, exportNamePdf, exportNamePng = scanExpCond(fileNames, rmBin, logPlot, avgPlot, font, IsItDLS = 0)      
    
    # Set export title if autoPlot is not selected
    else:
        if rmBin:
            endNamePdf = plotFileName + '_rmBin.pdf'
            endNamePng = plotFileName + '_rmBin.png'
        else:
            endNamePdf = plotFileName + '.pdf'
            endNamePng = plotFileName + '.png'

        if logPlot:
            exportNamePdf = os.path.join('02_MDI', '01_Overview', '02_Log', '01_PDF', endNamePdf)
            exportNamePng = os.path.join('02_MDI', '01_Overview', '02_Log', '02_PNG',endNamePng)
        else:
            exportNamePdf = os.path.join('02_MDI', '01_Overview', '01_Linear', '01_PDF', endNamePdf)
            exportNamePng = os.path.join('02_MDI', '01_Overview', '01_Linear', '02_PNG',endNamePng)

    if showTitle:                                 
        # Set overall title
        fig.suptitle(plotTitle, y=1.06)
        
    # Save plot    
    plt.savefig(os.path.join('02_Plots', exportNamePdf), bbox_inches='tight')
    plt.savefig(os.path.join('02_Plots', exportNamePng), bbox_inches='tight')
    
    # Don't show plot in jupyter notebook
    plt.close()
        
    return


def scanExpCond(fileNames, rmFirstBin, logPlot, avgPlot, font, IsItDLS):
        
        plotTitle = ''
        offset = 0
        
        # Extract data from filename
        fileName = fileNames[0]
        sampleID = fileName[0:3]
        sampleConc = fileName[4:8]

        if sampleID == 'H2O':
            plotTitle += 'Milli-Q Water'
        elif sampleID == 'JJ1':
            plotTitle += 'JJmAb 1'
        elif sampleID == 'JJ2':
            plotTitle += 'JJmAb 2'
        elif sampleID == 'JJ3':
            plotTitle += 'JJmAb 3'
        elif sampleID == 'JJ4':
            plotTitle += 'JJmAb 4'        
        
        if sampleConc != '00mg':
            plotTitle += ' ({:.0f} mg/mL), '.format(float(sampleConc[:2]))
        else:
            plotTitle += ', '
            
        if fileName[9:12] in set(['PS2', 'PS8', 'PO1', 'BRJ', 'AIR']):
            sampleID2 = fileName[9:12]
            sampleConc2 = fileName[13:18]
            offset = 10
                        
            if sampleID2 == 'PS2':
                plotTitle += 'Polysorbate 20 '
            elif sampleID2 == 'PS8':
                plotTitle += 'Polysorbate 80 '
            elif sampleID2 == 'PO1':
                plotTitle += 'Poloxamer 188 '
            elif sampleID2 == 'BRJ':
                plotTitle += 'Breach '
            elif sampleID2 == 'AIR':
                plotTitle += 'Air '
            
            if sampleConc2 != '000ug' and sampleConc2[-2:] != '%v':
                plotTitle += '({:.2f} mg/mL), '.format(float(sampleConc2[:3])/1000)
            elif  sampleConc2[-2:] == '%v':
                if font == 'LaTeX':
                    plotTitle += '({:.2f} \% vol.), '.format(float(sampleConc2[:3]))
                else:
                    plotTitle += '({:.2f} % vol.), '.format(float(sampleConc2[:3]))
            else: 
                plotTitle += ', '
         
        rotFreq = fileName[(9+offset):(15+offset)]
        expTime = fileName[(16+offset):(24+offset)]
        expMat = fileName[(25+offset):(28+offset)]
        
        plotTitle += '{:.0f} rpm'.format(float(rotFreq[:3]))
        
        plotTitle += ', {:.0f} min'.format(float(expTime[:2]))
            
        if expMat == 'PM':
            plotTitle += ', PMMA'
        elif expMat == 'ST':
            plotTitle += ', Steel'
        elif expMat == 'PA':
            plotTitle += ', PMMA + Air'

            
        # Set special name for SST measurement
        if sampleID == 'SST':
            plotTitle = 'NIST Standard (Polystyrene Beads of 200 nm Diameter)'
            
        # Set export title
        if avgPlot:
            if rmFirstBin:
                endNamePdf = fileNames[0] + '_rmBin_Avg.pdf'
                endNamePng = fileNames[0] + '_rmBin_Avg.png'
            else:
                endNamePdf = fileNames[0] + '_Avg.pdf'
                endNamePng = fileNames[0] + '_Avg.png'

            if IsItDLS:
                exportNamePdf = os.path.join('01_DLS', '02_Average','01_PDF', endNamePdf)
                exportNamePng = os.path.join('01_DLS', '02_Average','02_PNG', endNamePng)
            else:
                if logPlot:
                    exportNamePdf = os.path.join('02_MDI', '02_Average', '02_Log', '01_PDF', endNamePdf)
                    exportNamePng = os.path.join('02_MDI', '02_Average', '02_Log', '02_PNG',endNamePng)
                else:
                    exportNamePdf = os.path.join('02_MDI', '02_Average', '01_Linear', '01_PDF', endNamePdf)
                    exportNamePng = os.path.join('02_MDI', '02_Average', '01_Linear', '02_PNG',endNamePng)

        else:
            if rmFirstBin:
                endNamePdf = fileNames[0] + '_rmBin_Coll.pdf'
                endNamePng = fileNames[0] + '_rmBin_Coll.png'
            else:
                endNamePdf = fileNames[0] + '_Coll.pdf'
                endNamePng = fileNames[0] + '_Coll.png'

            if IsItDLS:
                exportNamePdf = os.path.join('01_DLS', '03_Collection','01_PDF', endNamePdf)
                exportNamePng = os.path.join('01_DLS', '03_Collection','02_PNG', endNamePng)
            else:
                if logPlot:
                    exportNamePdf = os.path.join('02_MDI', '03_Collection', '02_Log', '01_PDF', endNamePdf)
                    exportNamePng = os.path.join('02_MDI', '03_Collection', '02_Log', '02_PNG',endNamePng)
                else:
                    exportNamePdf = os.path.join('02_MDI', '03_Collection', '01_Linear', '01_PDF', endNamePdf)
                    exportNamePng = os.path.join('02_MDI', '03_Collection', '01_Linear', '02_PNG',endNamePng)

        return plotTitle, exportNamePdf, exportNamePng



def plotSurfPub(corrData, intData, countData, plotName, cmap, font, mat):

    """
    Plot MDI and DLS results in a subplot matrix
    """

    ######## MDI Plot ########    

    expNames = ('JJ1 \n Sheared', 'JJ1 \n Unsheared', '0.01 \n mg/mL', '0.05 \n mg/mL', '0.10 \n mg/mL', '0.25 \n mg/mL', 
                '0.01 \n mg/mL', '0.05 \n mg/mL', '0.10 \n mg/mL', '0.25 \n mg/mL', 
                '0.01 \n mg/mL', '0.05 \n mg/mL', '0.10 \n mg/mL', '0.25 \n mg/mL')

    fileNames = ['JJ1_01mg_200rpm_10-00min', 'JJ1_01mg_000rpm_10-00min',
                'JJ1_01mg_PO1_010ug_200rpm_10-00min', 
                'JJ1_01mg_PO1_050ug_200rpm_10-00min', 'JJ1_01mg_PO1_100ug_200rpm_10-00min', 
                'JJ1_01mg_PO1_250ug_200rpm_10-00min', 'JJ1_01mg_PS8_010ug_200rpm_10-00min', 
                'JJ1_01mg_PS8_050ug_200rpm_10-00min', 'JJ1_01mg_PS8_100ug_200rpm_10-00min', 
                'JJ1_01mg_PS8_250ug_200rpm_10-00min', 'JJ1_01mg_PS2_010ug_200rpm_10-00min', 
                'JJ1_01mg_PS2_050ug_200rpm_10-00min', 'JJ1_01mg_PS2_100ug_200rpm_10-00min', 
                'JJ1_01mg_PS2_250ug_200rpm_10-00min']
    
    fileNames = [fileName + '_' + mat for fileName in fileNames]
    fileNames[1] = 'JJ1_01mg_000rpm_10-00min_PM'


    # Set parameter lists containing plotting parameters
    # param = [autoPlot, rmFirstBin, logPlot, plotAverage, plotLimits, expNames, plotFileName, 
    #          plotTitle, colors, cmap, font, showTitle]
    param = [0, 0, 0, 1, (0, 4e6), expNames, plotName, [0], cmap, font]

    [autoPlot, rmFirstBin, logPlot, avgPlot, plotLimits, expNames, plotFileName, colors, cmap, font] = param

    # Preliminary parameters
    countLimY = plotLimits

    # Tuple containing experiment number strings
    expTuple = expNames

    numFiles = len(fileNames)

    # Initialize empty count dictionary (needed for grouped barplot)
    keyList = ['Total ($\geq$ 1 $\mu$m)', '1-2 $\mu$m', '2-5 $\mu$m', '5-10 $\mu$m']#, '10-25 $\mu$m']
    keyLength = len(keyList)
    keyContent = []
    meanCountDict = dict.fromkeys(keyList, keyContent)
    stdCountDict = dict.fromkeys(keyList, keyContent)

    # Adjust colormap settings
    #plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,6)))
    darkGray = [0.4, 0.4, 0.4, 1]   # RGBA value, in closed interval [0,1]
    colMap = getattr(plt.cm, cmap)(np.linspace(0,1,6))
    plt.rcParams["axes.prop_cycle"] = plt.cycler("color", np.insert(colMap, 0, darkGray, axis=0))

    countDataProc = {}
    countDataDel = {}

    firstTime = True

    # Loop over all given files
    for i, file in enumerate(fileNames):

        countDataProc[file] = np.insert(countData[file], 0, np.sum(countData[file], axis=0), axis=0)

        # Find max. count in given file
        countMax = np.max(countDataProc[file])

        # Check shape of correlation data array
        countShape = np.shape(countDataProc[file])

        # Calculate number of experiments (one duplicate per experiment)
        countExpNum = int(countShape[1] / 2)

        if firstTime:
            meanExpCount = np.zeros((countShape[0], countExpNum*numFiles))
            stdExpCount = np.zeros((countShape[0], countExpNum*numFiles)) 

        for j in range(int(countExpNum)):
            meanExpCount[:,i*countExpNum+j] = np.mean(countDataProc[file][:, (2*j):(2*j+2)], axis=1)
            stdExpCount[:,i*countExpNum+j] = np.std(countDataProc[file][:, (2*j):(2*j+2)], axis=1)#, ddof=1)

        if avgPlot:

            if firstTime:
                meanCount = np.zeros((countShape[0], numFiles))
                stdCount = np.zeros((countShape[0], numFiles))

            if i == 0:
                meanCount[:,i] = np.mean(meanExpCount[:,:(i*countExpNum+countExpNum)], axis=1)
                stdCount[:,i] = np.std(meanExpCount[:,:(i*countExpNum+countExpNum)], axis=1)#, ddof=1)
            else: 
                meanCount[:,i] = np.mean(meanExpCount[:,((i-1)*countExpNum+countExpNum):(i*countExpNum+countExpNum)], axis=1)
                stdCount[:,i] = np.std(meanExpCount[:,((i-1)*countExpNum+countExpNum):(i*countExpNum+countExpNum)], axis=1)#, ddof=1)

        if firstTime:

            figWidth = 20
            figHeight = figWidth/1.8
            
            fig = plt.figure(figsize=(figWidth, figHeight))
            gs = fig.add_gridspec(2,5, hspace=0.3, width_ratios=(2, 3, 3, 3, 0.1), height_ratios=(1.5, 1))
            ax = fig.add_subplot(gs[0, :])
            ax2 = fig.add_subplot(gs[1, 1])
            ax3 = fig.add_subplot(gs[1, 2])
            ax4 = fig.add_subplot(gs[1, 3])


            # Add some text for labels, title and custom x-axis tick labels, etc.
            #ax.set_xlabel('Experiments')
            ax.set_ylabel('Number Density [mL$^{-1}$]')
            ax.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                            bottom=True, top=False, left=True, right=True, direction='in')

            firstTime = False


    x = np.arange(len(expTuple))
    width = 0.2  # the width of the bars
    multiplier = 1

    for j in range(keyLength):                
        meanCountDict[keyList[j]] = meanCount[j,:]
        stdCountDict[keyList[j]] = stdCount[j,:]
        

    for attribute, measurement in meanCountDict.items():
        offset = width * multiplier
        rects = ax.bar(x + offset, measurement, width, label=attribute, alpha=0.7)
        ax.errorbar(x + offset, measurement, yerr=stdCountDict[attribute], ecolor='black',
                capsize=5, linestyle='')
        multiplier += 1
            

    # Set additional plotting parameters
    ax.set_xticks(x + (keyLength+1)/2*width, expTuple)
    ax.legend(ncols=1)

    if logPlot:
        ax.set_yscale('log')
    else:
        ax.set_yscale('linear')

    if plotLimits and logPlot:
        ax.set_ylim(countLimY)
    elif plotLimits:
        ax.set_ylim(countLimY)
        ax.ticklabel_format(style='sci', axis='y', scilimits=(0,0), useMathText=True)
    elif logPlot:
        ax.set_ylim(bottom=10)
    else:
        ax.set_ylim(bottom=0)
        ax.ticklabel_format(style='sci', axis='y', scilimits=(0,0), useMathText=True)

    ######## DLS Plots ########    

    plotLegend = ['Sheared', '0.01 mg/mL', '0.05 mg/mL', '0.1 mg/mL', '0.25 mg/mL', 'Unsheared']
    fileNames = ['JJ1_01mg_200rpm_10-00min', 'JJ1_01mg_PO1_010ug_200rpm_10-00min',
                'JJ1_01mg_PO1_050ug_200rpm_10-00min', 'JJ1_01mg_PO1_100ug_200rpm_10-00min', 
                'JJ1_01mg_PO1_250ug_200rpm_10-00min']
    
    fileNames = [fileName + '_' + mat for fileName in fileNames]

    fileNames.append('JJ1_01mg_000rpm_10-00min_PM')


    # Set parameter lists containing plotting parameters
    # param = [autoPlot, plotIntensity, plotAverage, plotNumberOfExp, plotFileName, 
    #          plotTitle, plotLegend, colors, cmap, font, showTitle]
    param = [0, 0, 1, 0, plotName, plotLegend, [0], cmap, font]

    [autoPlot, plotIntData, avgPlot, plotNumberOfExp, plotFileName, plotLegend, colors, cmap, font] = param

    # Preliminary parameters
    corrLimX = [0.5, 1e5]
    corrLimY = [-0.05, 1.05]
    intLimX = [0.4e-3, 10]
    intLimY = [-2,27]

    corrTextCoord = [2*corrLimX[0], 0.92*corrLimY[1]]
    intTextCoord = [2*intLimX[0], 0.92*intLimY[1]] 

    text = 'Number of Experiments: {:.0f}'

    multipleFiles = False

    # First time parameter (true only for first loop iteration)
    firstTime = True

    # Array to save whether intensity data is stored in each given file
    intMeasured = np.zeros(len(fileNames))

    # Preliminary check for intensity data
    for file, i in zip(fileNames, range(len(fileNames))):

        if intData[file].any():
            # Set bool to true
            intMeasured[i] = 1

    if len(fileNames) > 1:
        multipleFiles = True

    # Loop over all given files
    for i, file in enumerate(fileNames):

        # Check shape of correlation data array
        corrShape = np.shape(corrData[file])

        # Calculate number of experiments (one triplicate per experiment)
        corrExpNum = int((corrShape[1] - 1)/3)

        # Adjust colormap settings
        if multipleFiles:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,len(fileNames))))
        else:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,corrExpNum)))

        # Extract correlation data triplets (and compute row-wise mean & std. dev.)
        lagTime = corrData[file][:, 0]

        meanExpCorr = np.empty((corrShape[0], int(corrExpNum)))
        stdExpCorr = np.empty((corrShape[0], int(corrExpNum)))

        for j in range(int(corrExpNum)): 
            meanExpCorr[:,j] = np.mean(corrData[file][:, (3*j+1):(3*j+4)], axis=1)
            stdExpCorr[:,j] = np.std(corrData[file][:, (3*j+1):(3*j+4)], axis=1, ddof=1)

        if avgPlot and corrExpNum > 1:
            meanCorr = np.mean(meanExpCorr, axis=1)
            stdCorr = np.std(meanExpCorr, axis=1)

        elif avgPlot and corrExpNum == 1:
            meanCorr = meanExpCorr[:,0]
            stdCorr = stdExpCorr[:,0]

        else:
            meanCorr = meanExpCorr
            stdCorr = stdExpCorr


        if firstTime:
            # Skip first color
            ax2.plot([],[], label='_nolegend_')
            ax2.fill_between([], [], [], label='_nolegend_')
            # Set axis labels and legends
            ax2.set_xlabel('Lag Time [$\mu$s]')
            ax2.set_ylabel('Correlation Coefficient [-]')
            ax2.set_xlim(corrLimX)
            ax2.set_ylim(corrLimY)
            ax2.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                                bottom=True, top=False, left=True, right=True, direction='in')
            if multipleFiles and plotNumberOfExp:
                ax2.text(corrTextCoord[0], corrTextCoord[1], text.format(len(fileNames)), va='center')
            elif plotNumberOfExp:
                ax2.text(corrTextCoord[0], corrTextCoord[1], text.format(corrExpNum), va='center')
                
            firstTime = False
            

        if avgPlot:
            # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
            if colors[0]:
                ax2.semilogx(lagTime, meanCorr, color=colors[i])
                ax2.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, color=colors[i], alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)
            else:
                ax2.semilogx(lagTime, meanCorr)
                ax2.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)

            ax2.xaxis.get_major_locator().set_params(numticks=99)
            ax2.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')


        else:
            for j in range(int(corrExpNum)): 
                # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
                if colors[0]:
                    ax2.semilogx(lagTime, meanCorr[:,j], color=colors[i])
                    ax2.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                    color=colors[i], alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                    antialiased=True)
                else:
                    if multipleFiles:
                        fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                        ax2.semilogx(lagTime, meanCorr[:,j], color=fileColor)
                        ax2.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        color=fileColor, alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)
                    else:
                        ax2.semilogx(lagTime, meanCorr[:,j])
                        ax2.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)

                ax2.xaxis.get_major_locator().set_params(numticks=99)
                ax2.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')


        if plotLegend and i == (len(fileNames) - 1):
            ax2.legend(plotLegend, bbox_to_anchor=(-0.25, 1.05))

            if not avgPlot:
                for i in range(len(fileNames)):
                    leg = ax2.get_legend()
                    fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                    leg.legendHandles[i].set_color(fileColor)

    plotLegend = ['Sheared', '0.01 mg/mL', '0.05 mg/mL', '0.1 mg/mL', '0.25 mg/mL', 'Unsheared']
    fileNames = ['JJ1_01mg_200rpm_10-00min', 'JJ1_01mg_PS8_010ug_200rpm_10-00min',
                'JJ1_01mg_PS8_050ug_200rpm_10-00min', 'JJ1_01mg_PS8_100ug_200rpm_10-00min', 
                'JJ1_01mg_PS8_250ug_200rpm_10-00min']
    
    fileNames = [fileName + '_' + mat for fileName in fileNames]

    fileNames.append('JJ1_01mg_000rpm_10-00min_PM')

    # Set parameter lists containing plotting parameters
    # param = [autoPlot, plotIntensity, plotAverage, plotNumberOfExp, plotFileName, 
    #          plotTitle, plotLegend, colors, cmap, font, showTitle]
    param = [0, 0, 1, 0, plotName, plotLegend, [0], cmap, font]

    [autoPlot, plotIntData, avgPlot, plotNumberOfExp, plotFileName, plotLegend, colors, cmap, font] = param

    # Preliminary parameters
    corrLimX = [0.5, 1e5]
    corrLimY = [-0.05, 1.05]
    intLimX = [0.4e-3, 10]
    intLimY = [-2,27]

    corrTextCoord = [2*corrLimX[0], 0.92*corrLimY[1]]
    intTextCoord = [2*intLimX[0], 0.92*intLimY[1]] 

    text = 'Number of Experiments: {:.0f}'

    multipleFiles = False

    # First time parameter (true only for first loop iteration)
    firstTime = True

    # Array to save whether intensity data is stored in each given file
    intMeasured = np.zeros(len(fileNames))

    # Preliminary check for intensity data
    for file, i in zip(fileNames, range(len(fileNames))):

        if intData[file].any():
            # Set bool to true
            intMeasured[i] = 1

    if len(fileNames) > 1:
        multipleFiles = True


    # Loop over all given files
    for i, file in enumerate(fileNames):

        # Check shape of correlation data array
        corrShape = np.shape(corrData[file])

        # Calculate number of experiments (one triplicate per experiment)
        corrExpNum = int((corrShape[1] - 1)/3)

        # Adjust colormap settings
        if multipleFiles:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,len(fileNames))))
        else:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,corrExpNum)))

        # Extract correlation data triplets (and compute row-wise mean & std. dev.)
        lagTime = corrData[file][:, 0]

        meanExpCorr = np.empty((corrShape[0], int(corrExpNum)))
        stdExpCorr = np.empty((corrShape[0], int(corrExpNum)))

        for j in range(int(corrExpNum)): 
            meanExpCorr[:,j] = np.mean(corrData[file][:, (3*j+1):(3*j+4)], axis=1)
            stdExpCorr[:,j] = np.std(corrData[file][:, (3*j+1):(3*j+4)], axis=1, ddof=1)

        if avgPlot and corrExpNum > 1:
            meanCorr = np.mean(meanExpCorr, axis=1)
            stdCorr = np.std(meanExpCorr, axis=1)

        elif avgPlot and corrExpNum == 1:
            meanCorr = meanExpCorr[:,0]
            stdCorr = stdExpCorr[:,0]

        else:
            meanCorr = meanExpCorr
            stdCorr = stdExpCorr


        if firstTime:
            
            # Skip first color
            ax3.plot([],[], label='_nolegend_')
            ax3.fill_between([], [], [], label='_nolegend_')

            # Set axis labels and legends
            ax3.set_xlabel('Lag Time [$\mu$s]')
            #ax3.set_ylabel('Correlation Coefficient [-]')
            ax3.set_xlim(corrLimX)
            ax3.set_ylim(corrLimY)
            ax3.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                                bottom=True, top=False, left=True, right=True, direction='in')
            if multipleFiles and plotNumberOfExp:
                ax3.text(corrTextCoord[0], corrTextCoord[1], text.format(len(fileNames)), va='center')
            elif plotNumberOfExp:
                ax3.text(corrTextCoord[0], corrTextCoord[1], text.format(corrExpNum), va='center')

            firstTime = False

        if avgPlot:
            # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
            if colors[0]:
                ax3.semilogx(lagTime, meanCorr, color=colors[i])
                ax3.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, color=colors[i], alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)
            else:
                ax3.semilogx(lagTime, meanCorr)
                ax3.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)

            ax3.xaxis.get_major_locator().set_params(numticks=99)
            ax3.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')


        else:
            for j in range(int(corrExpNum)): 
                # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
                if colors[0]:
                    ax3.semilogx(lagTime, meanCorr[:,j], color=colors[i])
                    ax3.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                    color=colors[i], alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                    antialiased=True)
                else:
                    if multipleFiles:
                        fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                        ax3.semilogx(lagTime, meanCorr[:,j], color=fileColor)
                        ax3.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        color=fileColor, alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)
                    else:
                        ax3.semilogx(lagTime, meanCorr[:,j])
                        ax3.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)

                ax3.xaxis.get_major_locator().set_params(numticks=99)
                ax3.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')


        #if plotLegend and i == (len(fileNames) - 1):
        #    ax3.legend(plotLegend)

            if not avgPlot:
                for i in range(len(fileNames)):
                    leg = ax3.get_legend()
                    fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                    leg.legendHandles[i].set_color(fileColor)

    plotLegend = ['Sheared', '0.01 mg/mL', '0.05 mg/mL', '0.1 mg/mL', '0.25 mg/mL', 'Unsheared']
    fileNames = ['JJ1_01mg_200rpm_10-00min', 'JJ1_01mg_PS2_010ug_200rpm_10-00min',
                'JJ1_01mg_PS2_050ug_200rpm_10-00min', 'JJ1_01mg_PS2_100ug_200rpm_10-00min', 
                'JJ1_01mg_PS2_250ug_200rpm_10-00min']
    
    fileNames = [fileName + '_' + mat for fileName in fileNames]

    fileNames.append('JJ1_01mg_000rpm_10-00min_PM')

    # Set parameter lists containing plotting parameters
    # param = [autoPlot, plotIntensity, plotAverage, plotNumberOfExp, plotFileName, 
    #          plotTitle, plotLegend, colors, cmap, font, showTitle]
    param = [0, 0, 1, 0, plotName, plotLegend, [0], cmap, font]

    [autoPlot, plotIntData, avgPlot, plotNumberOfExp, plotFileName, plotLegend, colors, cmap, font] = param

    # Preliminary parameters
    corrLimX = [0.5, 1e5]
    corrLimY = [-0.05, 1.05]
    intLimX = [0.4e-3, 10]
    intLimY = [-2,27]

    corrTextCoord = [2*corrLimX[0], 0.92*corrLimY[1]]
    intTextCoord = [2*intLimX[0], 0.92*intLimY[1]] 

    text = 'Number of Experiments: {:.0f}'

    multipleFiles = False

    # First time parameter (true only for first loop iteration)
    firstTime = True

    # Array to save whether intensity data is stored in each given file
    intMeasured = np.zeros(len(fileNames))

    # Preliminary check for intensity data
    for file, i in zip(fileNames, range(len(fileNames))):

        if intData[file].any():
            # Set bool to true
            intMeasured[i] = 1

    if len(fileNames) > 1:
        multipleFiles = True


    # Loop over all given files
    for i, file in enumerate(fileNames):

        # Check shape of correlation data array
        corrShape = np.shape(corrData[file])

        # Calculate number of experiments (one triplicate per experiment)
        corrExpNum = int((corrShape[1] - 1)/3)

        # Adjust colormap settings
        if multipleFiles:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,len(fileNames))))
        else:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,corrExpNum)))

        # Extract correlation data triplets (and compute row-wise mean & std. dev.)
        lagTime = corrData[file][:, 0]

        meanExpCorr = np.empty((corrShape[0], int(corrExpNum)))
        stdExpCorr = np.empty((corrShape[0], int(corrExpNum)))

        for j in range(int(corrExpNum)): 
            meanExpCorr[:,j] = np.mean(corrData[file][:, (3*j+1):(3*j+4)], axis=1)
            stdExpCorr[:,j] = np.std(corrData[file][:, (3*j+1):(3*j+4)], axis=1, ddof=1)

        if avgPlot and corrExpNum > 1:
            meanCorr = np.mean(meanExpCorr, axis=1)
            stdCorr = np.std(meanExpCorr, axis=1)

        elif avgPlot and corrExpNum == 1:
            meanCorr = meanExpCorr[:,0]
            stdCorr = stdExpCorr[:,0]

        else:
            meanCorr = meanExpCorr
            stdCorr = stdExpCorr


        if firstTime:
            
            # Skip first color
            ax4.plot([],[], label='_nolegend_')
            ax4.fill_between([], [], [], label='_nolegend_')

            # Set axis labels and legends
            ax4.set_xlabel('Lag Time [$\mu$s]')
            #ax4.set_ylabel('Correlation Coefficient [-]')
            ax4.set_xlim(corrLimX)
            ax4.set_ylim(corrLimY)
            ax4.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                                bottom=True, top=False, left=True, right=True, direction='in')
            if multipleFiles and plotNumberOfExp:
                ax4.text(corrTextCoord[0], corrTextCoord[1], text.format(len(fileNames)), va='center')
            elif plotNumberOfExp:
                ax4.text(corrTextCoord[0], corrTextCoord[1], text.format(corrExpNum), va='center')

            firstTime = False

        if avgPlot:
            # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
            if colors[0]:
                ax4.semilogx(lagTime, meanCorr, color=colors[i])
                ax4.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, color=colors[i], alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)
            else:
                ax4.semilogx(lagTime, meanCorr)
                ax4.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)

            ax4.xaxis.get_major_locator().set_params(numticks=99)
            ax4.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')


        else:
            for j in range(int(corrExpNum)): 
                # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
                if colors[0]:
                    ax4.semilogx(lagTime, meanCorr[:,j], color=colors[i])
                    ax4.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                    color=colors[i], alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                    antialiased=True)
                else:
                    if multipleFiles:
                        fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                        ax4.semilogx(lagTime, meanCorr[:,j], color=fileColor)
                        ax4.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        color=fileColor, alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)
                    else:
                        ax4.semilogx(lagTime, meanCorr[:,j])
                        ax4.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)

                ax4.xaxis.get_major_locator().set_params(numticks=99)
                ax4.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')


        #if plotLegend and i == (len(fileNames) - 1):
        #    ax4.legend(plotLegend, bbox_to_anchor=(1, 1.05))

            if not avgPlot:
                for i in range(len(fileNames)):
                    leg = ax4.get_legend()
                    fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                    leg.legendHandles[i].set_color(fileColor)
                    
    #endNamePdf = plotFileName + '.pdf'
    #endNamePng = plotFileName + '.png'

    #if logPlot:
    #    exportNamePdf = os.path.join('MDI', 'Overview', 'Log', 'PDF', endNamePdf)
    #    exportNamePng = os.path.join('MDI', 'Overview', 'Log', 'PNG',endNamePng)
    #else:
    #    exportNamePdf = os.path.join('MDI', 'Overview', 'Linear', 'PDF', endNamePdf)
    #    exportNamePng = os.path.join('MDI', 'Overview', 'Linear', 'PNG',endNamePng)

    #if showTitle:                                 
        # Set overall title
        #fig.suptitle(plotTitle, y=1.06)

    # Save plot
    if mat == 'PM':   
        plt.savefig(os.path.join('02_Plots', '00_Summary', '01_PDF', 'Surf_PM.pdf'), bbox_inches='tight')
        plt.savefig(os.path.join('02_Plots', '00_Summary', '02_PNG', 'Surf_PM.png'), bbox_inches='tight')
    else:
        plt.savefig(os.path.join('02_Plots', '00_Summary', '01_PDF', 'Surf_ST.pdf'), bbox_inches='tight')
        plt.savefig(os.path.join('02_Plots', '00_Summary', '02_PNG', 'Surf_ST.png'), bbox_inches='tight')
    # Don't show plot in jupyter notebook
    #plt.close()


def plotConcPub(corrData, intData, countData, plotName, cmap, font, mat):

    """
    Plot MDI and DLS results in a subplot matrix
    """

    ######## MDI Plot ########    

    expNames = ('Sheared', 'Sheared, PS 20', 'Unsheared', 
                'Sheared', 'Sheared, PS 20', 'Unsheared', 
                'Sheared', 'Sheared, PS 20', 'Unsheared')

    fileNames = ['JJ1_01mg_200rpm_10-00min', 'JJ1_01mg_PS2_250ug_200rpm_10-00min', 'JJ1_01mg_000rpm_00-00min_NA',
                 'JJ1_10mg_200rpm_10-00min', 'JJ1_10mg_PS2_250ug_200rpm_10-00min', 'JJ1_10mg_000rpm_00-00min_NA',
                 'JJ1_50mg_200rpm_10-00min', 'JJ1_50mg_PS2_500ug_200rpm_10-00min', 'JJ1_50mg_000rpm_00-00min_NA']
    
    fileNames = [fileName + '_' + mat for fileName in fileNames]
    fileNames[2] = 'JJ1_01mg_000rpm_00-00min_NA'
    fileNames[5] = 'JJ1_10mg_000rpm_00-00min_NA'
    fileNames[8] = 'JJ1_50mg_000rpm_00-00min_NA'


    # Set parameter lists containing plotting parameters
    # param = [autoPlot, rmFirstBin, logPlot, plotAverage, plotLimits, expNames, plotFileName, 
    #          plotTitle, colors, cmap, font, showTitle]
    param = [0, 0, 0, 1, (0, 4e6), expNames, plotName, [0], cmap, font]

    [autoPlot, rmFirstBin, logPlot, avgPlot, plotLimits, expNames, plotFileName, colors, cmap, font] = param

    # Preliminary parameters
    countLimY = plotLimits

    # Tuple containing experiment number strings
    expTuple = expNames

    numFiles = len(fileNames)

    # Initialize empty count dictionary (needed for grouped barplot)
    keyList = ['Total ($\geq$ 1 $\mu$m)', '1-2 $\mu$m', '2-5 $\mu$m', '5-10 $\mu$m']#, '10-25 $\mu$m']
    keyLength = len(keyList)
    keyContent = []
    meanCountDict = dict.fromkeys(keyList, keyContent)
    stdCountDict = dict.fromkeys(keyList, keyContent)

    # Adjust colormap settings
    #plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,6)))
    darkGray = [0.4, 0.4, 0.4, 1]   # RGBA value, in closed interval [0,1]
    colMap = getattr(plt.cm, cmap)(np.linspace(0,1,6))
    plt.rcParams["axes.prop_cycle"] = plt.cycler("color", np.insert(colMap, 0, darkGray, axis=0))

    countDataProc = {}
    countDataDel = {}

    firstTime = True

    # Loop over all given files
    for i, file in enumerate(fileNames):

        countDataProc[file] = np.insert(countData[file], 0, np.sum(countData[file], axis=0), axis=0)

        # Find max. count in given file
        countMax = np.max(countDataProc[file])

        # Check shape of correlation data array
        countShape = np.shape(countDataProc[file])

        # Calculate number of experiments (one duplicate per experiment)
        countExpNum = int(countShape[1] / 2)

        if firstTime:
            meanExpCount = np.zeros((countShape[0], countExpNum*numFiles))
            stdExpCount = np.zeros((countShape[0], countExpNum*numFiles)) 

        for j in range(int(countExpNum)):
            meanExpCount[:,i*countExpNum+j] = np.mean(countDataProc[file][:, (2*j):(2*j+2)], axis=1)
            stdExpCount[:,i*countExpNum+j] = np.std(countDataProc[file][:, (2*j):(2*j+2)], axis=1)#, ddof=1)

        if avgPlot:

            if firstTime:
                meanCount = np.zeros((countShape[0], numFiles))
                stdCount = np.zeros((countShape[0], numFiles))

            if i == 0:
                meanCount[:,i] = np.mean(meanExpCount[:,:(i*countExpNum+countExpNum)], axis=1)
                stdCount[:,i] = np.std(meanExpCount[:,:(i*countExpNum+countExpNum)], axis=1)#, ddof=1)
            else: 
                meanCount[:,i] = np.mean(meanExpCount[:,((i-1)*countExpNum+countExpNum):(i*countExpNum+countExpNum)], axis=1)
                stdCount[:,i] = np.std(meanExpCount[:,((i-1)*countExpNum+countExpNum):(i*countExpNum+countExpNum)], axis=1)#, ddof=1)

        if firstTime:

            figWidth = 20
            figHeight = figWidth/1.8
            
            fig = plt.figure(figsize=(figWidth, figHeight))
            gs = fig.add_gridspec(2,5, hspace=0.3, width_ratios=(0.5, 3, 3, 3, 0.1), height_ratios=(1.5, 1))
            ax = fig.add_subplot(gs[0, :])
            ax2 = fig.add_subplot(gs[1, 1])
            ax3 = fig.add_subplot(gs[1, 2])
            ax4 = fig.add_subplot(gs[1, 3])


            # Add some text for labels, title and custom x-axis tick labels, etc.
            #ax.set_xlabel('Experiments')
            ax.set_ylabel('Number Density [mL$^{-1}$]')
            ax.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                            bottom=True, top=False, left=True, right=True, direction='in')

            firstTime = False


    x = np.arange(len(expTuple))
    width = 0.2  # the width of the bars
    multiplier = 1

    for j in range(keyLength):                
        meanCountDict[keyList[j]] = meanCount[j,:]
        stdCountDict[keyList[j]] = stdCount[j,:]
        

    for attribute, measurement in meanCountDict.items():
        offset = width * multiplier
        rects = ax.bar(x + offset, measurement, width, label=attribute, alpha=0.7)
        ax.errorbar(x + offset, measurement, yerr=stdCountDict[attribute], ecolor='black',
                capsize=5, linestyle='')
        multiplier += 1
            

    # Set additional plotting parameters
    ax.set_xticks(x + (keyLength+1)/2*width, expTuple)
    ax.legend(ncols=1)

    if logPlot:
        ax.set_yscale('log')
    else:
        ax.set_yscale('linear')

    if plotLimits and logPlot:
        ax.set_ylim(countLimY)
    elif plotLimits:
        ax.set_ylim(countLimY)
        ax.ticklabel_format(style='sci', axis='y', scilimits=(0,0), useMathText=True)
    elif logPlot:
        ax.set_ylim(bottom=10)
    else:
        ax.set_ylim(bottom=0)
        ax.ticklabel_format(style='sci', axis='y', scilimits=(0,0), useMathText=True)

    ######## DLS Plots ########    

    plotLegend = ['Sheared', 'Sheared, PS 20', 'Unsheared']
    fileNames = ['JJ1_01mg_200rpm_10-00min', 'JJ1_01mg_PS2_250ug_200rpm_10-00min', 'JJ1_01mg_000rpm_00-00min_NA']

    fileNames = [fileName + '_' + mat for fileName in fileNames]

    fileNames[2] = 'JJ1_01mg_000rpm_00-00min_NA'


    # Set parameter lists containing plotting parameters
    # param = [autoPlot, plotIntensity, plotAverage, plotNumberOfExp, plotFileName, 
    #          plotTitle, plotLegend, colors, cmap, font, showTitle]
    param = [0, 0, 1, 0, plotName, plotLegend, [0], cmap, font]

    [autoPlot, plotIntData, avgPlot, plotNumberOfExp, plotFileName, plotLegend, colors, cmap, font] = param

    # Preliminary parameters
    corrLimX = [0.5, 1e5]
    corrLimY = [-0.05, 1.05]
    intLimX = [0.4e-3, 10]
    intLimY = [-2,27]

    corrTextCoord = [2*corrLimX[0], 0.92*corrLimY[1]]
    intTextCoord = [2*intLimX[0], 0.92*intLimY[1]] 

    text = 'Number of Experiments: {:.0f}'

    multipleFiles = False

    # First time parameter (true only for first loop iteration)
    firstTime = True

    # Array to save whether intensity data is stored in each given file
    intMeasured = np.zeros(len(fileNames))

    # Preliminary check for intensity data
    for file, i in zip(fileNames, range(len(fileNames))):

        if intData[file].any():
            # Set bool to true
            intMeasured[i] = 1

    if len(fileNames) > 1:
        multipleFiles = True

    # Loop over all given files
    for i, file in enumerate(fileNames):

        # Check shape of correlation data array
        corrShape = np.shape(corrData[file])

        # Calculate number of experiments (one triplicate per experiment)
        corrExpNum = int((corrShape[1] - 1)/3)

        # Adjust colormap settings
        if multipleFiles:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,len(fileNames))))
        else:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,corrExpNum)))

        # Extract correlation data triplets (and compute row-wise mean & std. dev.)
        lagTime = corrData[file][:, 0]

        meanExpCorr = np.empty((corrShape[0], int(corrExpNum)))
        stdExpCorr = np.empty((corrShape[0], int(corrExpNum)))

        for j in range(int(corrExpNum)): 
            meanExpCorr[:,j] = np.mean(corrData[file][:, (3*j+1):(3*j+4)], axis=1)
            stdExpCorr[:,j] = np.std(corrData[file][:, (3*j+1):(3*j+4)], axis=1, ddof=1)

        if avgPlot and corrExpNum > 1:
            meanCorr = np.mean(meanExpCorr, axis=1)
            stdCorr = np.std(meanExpCorr, axis=1)

        elif avgPlot and corrExpNum == 1:
            meanCorr = meanExpCorr[:,0]
            stdCorr = stdExpCorr[:,0]

        else:
            meanCorr = meanExpCorr
            stdCorr = stdExpCorr


        if firstTime:
            # Skip first color
            ax2.plot([],[], label='_nolegend_')
            ax2.fill_between([], [], [], label='_nolegend_')
            # Set axis labels and legends
            ax2.set_xlabel('Lag Time [$\mu$s]')
            ax2.set_ylabel('Correlation Coefficient [-]')
            ax2.set_xlim(corrLimX)
            ax2.set_ylim(corrLimY)
            ax2.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                                bottom=True, top=False, left=True, right=True, direction='in')
            if multipleFiles and plotNumberOfExp:
                ax2.text(corrTextCoord[0], corrTextCoord[1], text.format(len(fileNames)), va='center')
            elif plotNumberOfExp:
                ax2.text(corrTextCoord[0], corrTextCoord[1], text.format(corrExpNum), va='center')
                
            firstTime = False
            

        if avgPlot:
            # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
            if colors[0]:
                ax2.semilogx(lagTime, meanCorr, color=colors[i])
                ax2.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, color=colors[i], alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)
            else:
                ax2.semilogx(lagTime, meanCorr)
                ax2.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)

            ax2.xaxis.get_major_locator().set_params(numticks=99)
            ax2.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')


        else:
            for j in range(int(corrExpNum)): 
                # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
                if colors[0]:
                    ax2.semilogx(lagTime, meanCorr[:,j], color=colors[i])
                    ax2.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                    color=colors[i], alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                    antialiased=True)
                else:
                    if multipleFiles:
                        fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                        ax2.semilogx(lagTime, meanCorr[:,j], color=fileColor)
                        ax2.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        color=fileColor, alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)
                    else:
                        ax2.semilogx(lagTime, meanCorr[:,j])
                        ax2.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)

                ax2.xaxis.get_major_locator().set_params(numticks=99)
                ax2.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')

        # Skip one color
        ax2.plot([],[], label='_nolegend_')
        ax2.fill_between([], [], [], label='_nolegend_')

        if plotLegend and i == (len(fileNames) - 1):
            ax2.legend(plotLegend, bbox_to_anchor=(-0.25, 1.05))

            if not avgPlot:
                for i in range(len(fileNames)):
                    leg = ax2.get_legend()
                    fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                    leg.legendHandles[i].set_color(fileColor)

    plotLegend = ['Sheared', 'Sheared, PS 20', 'Unsheared']
    fileNames = ['JJ1_10mg_200rpm_10-00min', 'JJ1_10mg_PS2_250ug_200rpm_10-00min', 'JJ1_10mg_000rpm_00-00min_NA']
    
    fileNames = [fileName + '_' + mat for fileName in fileNames]

    fileNames[2] = 'JJ1_10mg_000rpm_00-00min_NA'

    # Set parameter lists containing plotting parameters
    # param = [autoPlot, plotIntensity, plotAverage, plotNumberOfExp, plotFileName, 
    #          plotTitle, plotLegend, colors, cmap, font, showTitle]
    param = [0, 0, 1, 0, plotName, plotLegend, [0], cmap, font]

    [autoPlot, plotIntData, avgPlot, plotNumberOfExp, plotFileName, plotLegend, colors, cmap, font] = param

    # Preliminary parameters
    corrLimX = [0.5, 1e5]
    corrLimY = [-0.05, 1.05]
    intLimX = [0.4e-3, 10]
    intLimY = [-2,27]

    corrTextCoord = [2*corrLimX[0], 0.92*corrLimY[1]]
    intTextCoord = [2*intLimX[0], 0.92*intLimY[1]] 

    text = 'Number of Experiments: {:.0f}'

    multipleFiles = False

    # First time parameter (true only for first loop iteration)
    firstTime = True

    # Array to save whether intensity data is stored in each given file
    intMeasured = np.zeros(len(fileNames))

    # Preliminary check for intensity data
    for file, i in zip(fileNames, range(len(fileNames))):

        if intData[file].any():
            # Set bool to true
            intMeasured[i] = 1

    if len(fileNames) > 1:
        multipleFiles = True


    # Loop over all given files
    for i, file in enumerate(fileNames):

        # Check shape of correlation data array
        corrShape = np.shape(corrData[file])

        # Calculate number of experiments (one triplicate per experiment)
        corrExpNum = int((corrShape[1] - 1)/3)

        # Adjust colormap settings
        if multipleFiles:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,len(fileNames))))
        else:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,corrExpNum)))

        # Extract correlation data triplets (and compute row-wise mean & std. dev.)
        lagTime = corrData[file][:, 0]

        meanExpCorr = np.empty((corrShape[0], int(corrExpNum)))
        stdExpCorr = np.empty((corrShape[0], int(corrExpNum)))

        for j in range(int(corrExpNum)): 
            meanExpCorr[:,j] = np.mean(corrData[file][:, (3*j+1):(3*j+4)], axis=1)
            stdExpCorr[:,j] = np.std(corrData[file][:, (3*j+1):(3*j+4)], axis=1, ddof=1)

        if avgPlot and corrExpNum > 1:
            meanCorr = np.mean(meanExpCorr, axis=1)
            stdCorr = np.std(meanExpCorr, axis=1)

        elif avgPlot and corrExpNum == 1:
            meanCorr = meanExpCorr[:,0]
            stdCorr = stdExpCorr[:,0]

        else:
            meanCorr = meanExpCorr
            stdCorr = stdExpCorr


        if firstTime:
            
            # Skip first color
            ax3.plot([],[], label='_nolegend_')
            ax3.fill_between([], [], [], label='_nolegend_')

            # Set axis labels and legends
            ax3.set_xlabel('Lag Time [$\mu$s]')
            #ax3.set_ylabel('Correlation Coefficient [-]')
            ax3.set_xlim(corrLimX)
            ax3.set_ylim(corrLimY)
            ax3.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                                bottom=True, top=False, left=True, right=True, direction='in')
            if multipleFiles and plotNumberOfExp:
                ax3.text(corrTextCoord[0], corrTextCoord[1], text.format(len(fileNames)), va='center')
            elif plotNumberOfExp:
                ax3.text(corrTextCoord[0], corrTextCoord[1], text.format(corrExpNum), va='center')

            firstTime = False

        if avgPlot:
            # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
            if colors[0]:
                ax3.semilogx(lagTime, meanCorr, color=colors[i])
                ax3.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, color=colors[i], alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)
            else:
                ax3.semilogx(lagTime, meanCorr)
                ax3.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)

            ax3.xaxis.get_major_locator().set_params(numticks=99)
            ax3.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')


        else:
            for j in range(int(corrExpNum)): 
                # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
                if colors[0]:
                    ax3.semilogx(lagTime, meanCorr[:,j], color=colors[i])
                    ax3.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                    color=colors[i], alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                    antialiased=True)
                else:
                    if multipleFiles:
                        fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                        ax3.semilogx(lagTime, meanCorr[:,j], color=fileColor)
                        ax3.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        color=fileColor, alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)
                    else:
                        ax3.semilogx(lagTime, meanCorr[:,j])
                        ax3.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)

                ax3.xaxis.get_major_locator().set_params(numticks=99)
                ax3.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')

        # Skip one color
        ax3.plot([],[], label='_nolegend_')
        ax3.fill_between([], [], [], label='_nolegend_')

        #if plotLegend and i == (len(fileNames) - 1):
        #    ax3.legend(plotLegend)

        if not avgPlot:
            for i in range(len(fileNames)):
                leg = ax3.get_legend()
                fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                leg.legendHandles[i].set_color(fileColor)

    plotLegend = ['Sheared', 'Sheared, PS 20', 'Unsheared']
    fileNames = ['JJ1_50mg_200rpm_10-00min', 'JJ1_50mg_PS2_500ug_200rpm_10-00min', 'JJ1_50mg_000rpm_00-00min_NA']
    
    fileNames = [fileName + '_' + mat for fileName in fileNames]

    fileNames[2] = 'JJ1_50mg_000rpm_00-00min_NA'

    # Set parameter lists containing plotting parameters
    # param = [autoPlot, plotIntensity, plotAverage, plotNumberOfExp, plotFileName, 
    #          plotTitle, plotLegend, colors, cmap, font, showTitle]
    param = [0, 0, 1, 0, plotName, plotLegend, [0], cmap, font]

    [autoPlot, plotIntData, avgPlot, plotNumberOfExp, plotFileName, plotLegend, colors, cmap, font] = param

    # Preliminary parameters
    corrLimX = [0.5, 1e5]
    corrLimY = [-0.05, 1.05]
    intLimX = [0.4e-3, 10]
    intLimY = [-2,27]

    corrTextCoord = [2*corrLimX[0], 0.92*corrLimY[1]]
    intTextCoord = [2*intLimX[0], 0.92*intLimY[1]] 

    text = 'Number of Experiments: {:.0f}'

    multipleFiles = False

    # First time parameter (true only for first loop iteration)
    firstTime = True

    # Array to save whether intensity data is stored in each given file
    intMeasured = np.zeros(len(fileNames))

    # Preliminary check for intensity data
    for file, i in zip(fileNames, range(len(fileNames))):

        if intData[file].any():
            # Set bool to true
            intMeasured[i] = 1

    if len(fileNames) > 1:
        multipleFiles = True


    # Loop over all given files
    for i, file in enumerate(fileNames):

        # Check shape of correlation data array
        corrShape = np.shape(corrData[file])

        # Calculate number of experiments (one triplicate per experiment)
        corrExpNum = int((corrShape[1] - 1)/3)

        # Adjust colormap settings
        if multipleFiles:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,len(fileNames))))
        else:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,corrExpNum)))

        # Extract correlation data triplets (and compute row-wise mean & std. dev.)
        lagTime = corrData[file][:, 0]

        meanExpCorr = np.empty((corrShape[0], int(corrExpNum)))
        stdExpCorr = np.empty((corrShape[0], int(corrExpNum)))

        for j in range(int(corrExpNum)): 
            meanExpCorr[:,j] = np.mean(corrData[file][:, (3*j+1):(3*j+4)], axis=1)
            stdExpCorr[:,j] = np.std(corrData[file][:, (3*j+1):(3*j+4)], axis=1, ddof=1)

        if avgPlot and corrExpNum > 1:
            meanCorr = np.mean(meanExpCorr, axis=1)
            stdCorr = np.std(meanExpCorr, axis=1)

        elif avgPlot and corrExpNum == 1:
            meanCorr = meanExpCorr[:,0]
            stdCorr = stdExpCorr[:,0]

        else:
            meanCorr = meanExpCorr
            stdCorr = stdExpCorr


        if firstTime:
            
            # Skip first color
            ax4.plot([],[], label='_nolegend_')
            ax4.fill_between([], [], [], label='_nolegend_')

            # Set axis labels and legends
            ax4.set_xlabel('Lag Time [$\mu$s]')
            #ax4.set_ylabel('Correlation Coefficient [-]')
            ax4.set_xlim(corrLimX)
            ax4.set_ylim(corrLimY)
            ax4.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                                bottom=True, top=False, left=True, right=True, direction='in')
            if multipleFiles and plotNumberOfExp:
                ax4.text(corrTextCoord[0], corrTextCoord[1], text.format(len(fileNames)), va='center')
            elif plotNumberOfExp:
                ax4.text(corrTextCoord[0], corrTextCoord[1], text.format(corrExpNum), va='center')

            firstTime = False

        if avgPlot:
            # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
            if colors[0]:
                ax4.semilogx(lagTime, meanCorr, color=colors[i])
                ax4.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, color=colors[i], alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)
            else:
                ax4.semilogx(lagTime, meanCorr)
                ax4.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)

            ax4.xaxis.get_major_locator().set_params(numticks=99)
            ax4.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')


        else:
            for j in range(int(corrExpNum)): 
                # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
                if colors[0]:
                    ax4.semilogx(lagTime, meanCorr[:,j], color=colors[i])
                    ax4.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                    color=colors[i], alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                    antialiased=True)
                else:
                    if multipleFiles:
                        fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                        ax4.semilogx(lagTime, meanCorr[:,j], color=fileColor)
                        ax4.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        color=fileColor, alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)
                    else:
                        ax4.semilogx(lagTime, meanCorr[:,j])
                        ax4.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)

                ax4.xaxis.get_major_locator().set_params(numticks=99)
                ax4.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')

        # Skip one color
        ax4.plot([],[], label='_nolegend_')
        ax4.fill_between([], [], [], label='_nolegend_')

        #if plotLegend and i == (len(fileNames) - 1):
        #    ax4.legend(plotLegend, bbox_to_anchor=(1, 1.05))

        if not avgPlot:
            for i in range(len(fileNames)):
                leg = ax4.get_legend()
                fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                leg.legendHandles[i].set_color(fileColor)
                    
    #endNamePdf = plotFileName + '.pdf'
    #endNamePng = plotFileName + '.png'

    #if logPlot:
    #    exportNamePdf = os.path.join('MDI', 'Overview', 'Log', 'PDF', endNamePdf)
    #    exportNamePng = os.path.join('MDI', 'Overview', 'Log', 'PNG',endNamePng)
    #else:
    #    exportNamePdf = os.path.join('MDI', 'Overview', 'Linear', 'PDF', endNamePdf)
    #    exportNamePng = os.path.join('MDI', 'Overview', 'Linear', 'PNG',endNamePng)

    #if showTitle:                                 
        # Set overall title
        #fig.suptitle(plotTitle, y=1.06)

    # Save plot    
    if mat == 'PM':   
        plt.savefig(os.path.join('02_Plots', '00_Summary', '01_PDF', 'Conc_PM.pdf'), bbox_inches='tight')
        plt.savefig(os.path.join('02_Plots', '00_Summary', '02_PNG', 'Conc_PM.png'), bbox_inches='tight')
    else:
        plt.savefig(os.path.join('02_Plots', '00_Summary', '01_PDF', 'Conc_ST.pdf'), bbox_inches='tight')
        plt.savefig(os.path.join('02_Plots', '00_Summary', '02_PNG', 'Conc_ST.png'), bbox_inches='tight')

    # Don't show plot in jupyter notebook
    #plt.close()

def plotAirPub(corrData, intData, countData, plotName, cmap, font, mat):

    """
    Plot MDI and DLS results in a subplot matrix
    """

    ######## MDI Plot ########    

    expNames = ('Sheared', 'Sheared, Air', 'Unsheared', 
                'Sheared', 'Sheared, Air', 'Unsheared', 
                'Sheared', 'Sheared, Air', 'Unsheared')

    fileNames = ['JJ1_01mg_200rpm_10-00min', 'JJ1_01mg_AIR_010%v_200rpm_10-00min', 'JJ1_01mg_000rpm_00-00min_NA',
                 'JJ1_10mg_200rpm_10-00min', 'JJ1_10mg_AIR_010%v_200rpm_10-00min', 'JJ1_10mg_000rpm_00-00min_NA',
                 'JJ1_50mg_200rpm_10-00min', 'JJ1_50mg_AIR_010%v_200rpm_10-00min', 'JJ1_50mg_000rpm_00-00min_NA']
    
    fileNames = [fileName + '_' + mat for fileName in fileNames]
    fileNames[2] = 'JJ1_01mg_000rpm_00-00min_NA'
    fileNames[5] = 'JJ1_10mg_000rpm_00-00min_NA'
    fileNames[8] = 'JJ1_50mg_000rpm_00-00min_NA'


    # Set parameter lists containing plotting parameters
    # param = [autoPlot, rmFirstBin, logPlot, plotAverage, plotLimits, expNames, plotFileName, 
    #          plotTitle, colors, cmap, font, showTitle]
    param = [0, 0, 0, 1, (10, 2e8), expNames, plotName, [0], cmap, font]

    [autoPlot, rmFirstBin, logPlot, avgPlot, plotLimits, expNames, plotFileName, colors, cmap, font] = param

    # Preliminary parameters
    countLimY = plotLimits

    # Tuple containing experiment number strings
    expTuple = expNames

    numFiles = len(fileNames)

    # Initialize empty count dictionary (needed for grouped barplot)
    keyList = ['Total ($\geq$ 1 $\mu$m)', '1-2 $\mu$m', '2-5 $\mu$m', '5-10 $\mu$m']#, '10-25 $\mu$m']
    keyLength = len(keyList)
    keyContent = []
    meanCountDict = dict.fromkeys(keyList, keyContent)
    stdCountDict = dict.fromkeys(keyList, keyContent)

    # Adjust colormap settings
    #plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,6)))
    darkGray = [0.4, 0.4, 0.4, 1]   # RGBA value, in closed interval [0,1]
    colMap = getattr(plt.cm, cmap)(np.linspace(0,1,6))
    plt.rcParams["axes.prop_cycle"] = plt.cycler("color", np.insert(colMap, 0, darkGray, axis=0))

    countDataProc = {}
    countDataDel = {}

    firstTime = True

    # Loop over all given files
    for i, file in enumerate(fileNames):

        countDataProc[file] = np.insert(countData[file], 0, np.sum(countData[file], axis=0), axis=0)

        # Find max. count in given file
        countMax = np.max(countDataProc[file])

        # Check shape of correlation data array
        countShape = np.shape(countDataProc[file])

        # Calculate number of experiments (one duplicate per experiment)
        countExpNum = int(countShape[1] / 2)

        if firstTime:
            meanExpCount = np.zeros((countShape[0], countExpNum*numFiles))
            stdExpCount = np.zeros((countShape[0], countExpNum*numFiles)) 

        for j in range(int(countExpNum)):
            meanExpCount[:,i*countExpNum+j] = np.mean(countDataProc[file][:, (2*j):(2*j+2)], axis=1)
            stdExpCount[:,i*countExpNum+j] = np.std(countDataProc[file][:, (2*j):(2*j+2)], axis=1)#, ddof=1)

        if avgPlot:

            if firstTime:
                meanCount = np.zeros((countShape[0], numFiles))
                stdCount = np.zeros((countShape[0], numFiles))

            if i == 0:
                meanCount[:,i] = np.mean(meanExpCount[:,:(i*countExpNum+countExpNum)], axis=1)
                stdCount[:,i] = np.std(meanExpCount[:,:(i*countExpNum+countExpNum)], axis=1)#, ddof=1)
            else: 
                meanCount[:,i] = np.mean(meanExpCount[:,((i-1)*countExpNum+countExpNum):(i*countExpNum+countExpNum)], axis=1)
                stdCount[:,i] = np.std(meanExpCount[:,((i-1)*countExpNum+countExpNum):(i*countExpNum+countExpNum)], axis=1)#, ddof=1)

        if firstTime:

            figWidth = 20
            figHeight = figWidth/1.8
            
            fig = plt.figure(figsize=(figWidth, figHeight))
            gs = fig.add_gridspec(2,5, hspace=0.3, width_ratios=(0.5, 3, 3, 3, 0.1), height_ratios=(1.5, 1))
            ax = fig.add_subplot(gs[0, :])
            ax2 = fig.add_subplot(gs[1, 1])
            ax3 = fig.add_subplot(gs[1, 2])
            ax4 = fig.add_subplot(gs[1, 3])


            # Add some text for labels, title and custom x-axis tick labels, etc.
            #ax.set_xlabel('Experiments')
            ax.set_ylabel('Number Density [mL$^{-1}$]')
            ax.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                            bottom=True, top=False, left=True, right=True, direction='in')

            firstTime = False


    x = np.arange(len(expTuple))
    width = 0.2  # the width of the bars
    multiplier = 1

    for j in range(keyLength):                
        meanCountDict[keyList[j]] = meanCount[j,:]
        stdCountDict[keyList[j]] = stdCount[j,:]
        

    for attribute, measurement in meanCountDict.items():
        offset = width * multiplier
        rects = ax.bar(x + offset, measurement, width, label=attribute, alpha=0.7)
        ax.errorbar(x + offset, measurement, yerr=stdCountDict[attribute], ecolor='black',
                capsize=5, linestyle='')
        multiplier += 1
            

    # Set additional plotting parameters
    ax.set_xticks(x + (keyLength+1)/2*width, expTuple)
    ax.legend(ncols=1)

    if logPlot:
        ax.set_yscale('log')
    else:
        ax.set_yscale('linear')

    if plotLimits and logPlot:
        ax.set_ylim(countLimY)
    elif plotLimits:
        ax.set_ylim(countLimY)
        ax.ticklabel_format(style='sci', axis='y', scilimits=(0,0), useMathText=True)
    elif logPlot:
        ax.set_ylim(bottom=10)
    else:
        ax.set_ylim(bottom=0)
        ax.ticklabel_format(style='sci', axis='y', scilimits=(0,0), useMathText=True)

    ######## DLS Plots ########    

    plotLegend = ['Sheared', 'Sheared, Air', 'Unsheared']
    fileNames = ['JJ1_01mg_200rpm_10-00min', 'JJ1_01mg_AIR_010%v_200rpm_10-00min', 'JJ1_01mg_000rpm_00-00min_NA']

    fileNames = [fileName + '_' + mat for fileName in fileNames]

    fileNames[2] = 'JJ1_01mg_000rpm_00-00min_NA'


    # Set parameter lists containing plotting parameters
    # param = [autoPlot, plotIntensity, plotAverage, plotNumberOfExp, plotFileName, 
    #          plotTitle, plotLegend, colors, cmap, font, showTitle]
    param = [0, 0, 1, 0, plotName, plotLegend, [0], cmap, font]

    [autoPlot, plotIntData, avgPlot, plotNumberOfExp, plotFileName, plotLegend, colors, cmap, font] = param

    # Preliminary parameters
    corrLimX = [0.5, 1e5]
    corrLimY = [-0.05, 1.05]
    intLimX = [0.4e-3, 10]
    intLimY = [-2,27]

    corrTextCoord = [2*corrLimX[0], 0.92*corrLimY[1]]
    intTextCoord = [2*intLimX[0], 0.92*intLimY[1]] 

    text = 'Number of Experiments: {:.0f}'

    multipleFiles = False

    # First time parameter (true only for first loop iteration)
    firstTime = True

    # Array to save whether intensity data is stored in each given file
    intMeasured = np.zeros(len(fileNames))

    # Preliminary check for intensity data
    for file, i in zip(fileNames, range(len(fileNames))):

        if intData[file].any():
            # Set bool to true
            intMeasured[i] = 1

    if len(fileNames) > 1:
        multipleFiles = True

    # Loop over all given files
    for i, file in enumerate(fileNames):

        # Check shape of correlation data array
        corrShape = np.shape(corrData[file])

        # Calculate number of experiments (one triplicate per experiment)
        corrExpNum = int((corrShape[1] - 1)/3)

        # Adjust colormap settings
        if multipleFiles:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,len(fileNames))))
        else:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,corrExpNum)))

        # Extract correlation data triplets (and compute row-wise mean & std. dev.)
        lagTime = corrData[file][:, 0]

        meanExpCorr = np.empty((corrShape[0], int(corrExpNum)))
        stdExpCorr = np.empty((corrShape[0], int(corrExpNum)))

        for j in range(int(corrExpNum)): 
            meanExpCorr[:,j] = np.mean(corrData[file][:, (3*j+1):(3*j+4)], axis=1)
            stdExpCorr[:,j] = np.std(corrData[file][:, (3*j+1):(3*j+4)], axis=1, ddof=1)

        if avgPlot and corrExpNum > 1:
            meanCorr = np.mean(meanExpCorr, axis=1)
            stdCorr = np.std(meanExpCorr, axis=1)

        elif avgPlot and corrExpNum == 1:
            meanCorr = meanExpCorr[:,0]
            stdCorr = stdExpCorr[:,0]

        else:
            meanCorr = meanExpCorr
            stdCorr = stdExpCorr


        if firstTime:
            # Skip first color
            ax2.plot([],[], label='_nolegend_')
            ax2.fill_between([], [], [], label='_nolegend_')
            # Set axis labels and legends
            ax2.set_xlabel('Lag Time [$\mu$s]')
            ax2.set_ylabel('Correlation Coefficient [-]')
            ax2.set_xlim(corrLimX)
            ax2.set_ylim(corrLimY)
            ax2.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                                bottom=True, top=False, left=True, right=True, direction='in')
            if multipleFiles and plotNumberOfExp:
                ax2.text(corrTextCoord[0], corrTextCoord[1], text.format(len(fileNames)), va='center')
            elif plotNumberOfExp:
                ax2.text(corrTextCoord[0], corrTextCoord[1], text.format(corrExpNum), va='center')
                
            firstTime = False
            

        if avgPlot:
            # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
            if colors[0]:
                ax2.semilogx(lagTime, meanCorr, color=colors[i])
                ax2.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, color=colors[i], alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)
            else:
                ax2.semilogx(lagTime, meanCorr)
                ax2.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)

            ax2.xaxis.get_major_locator().set_params(numticks=99)
            ax2.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')


        else:
            for j in range(int(corrExpNum)): 
                # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
                if colors[0]:
                    ax2.semilogx(lagTime, meanCorr[:,j], color=colors[i])
                    ax2.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                    color=colors[i], alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                    antialiased=True)
                else:
                    if multipleFiles:
                        fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                        ax2.semilogx(lagTime, meanCorr[:,j], color=fileColor)
                        ax2.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        color=fileColor, alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)
                    else:
                        ax2.semilogx(lagTime, meanCorr[:,j])
                        ax2.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)

                ax2.xaxis.get_major_locator().set_params(numticks=99)
                ax2.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')

        # Skip one color
        ax2.plot([],[], label='_nolegend_')
        ax2.fill_between([], [], [], label='_nolegend_')

        if plotLegend and i == (len(fileNames) - 1):
            ax2.legend(plotLegend, bbox_to_anchor=(-0.25, 1.05))

            if not avgPlot:
                for i in range(len(fileNames)):
                    leg = ax2.get_legend()
                    fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                    leg.legendHandles[i].set_color(fileColor)

    plotLegend = ['Sheared', 'Sheared, Air', 'Unsheared']
    fileNames = ['JJ1_10mg_200rpm_10-00min', 'JJ1_10mg_AIR_010%v_200rpm_10-00min', 'JJ1_10mg_000rpm_00-00min_NA']
    
    fileNames = [fileName + '_' + mat for fileName in fileNames]

    fileNames[2] = 'JJ1_10mg_000rpm_00-00min_NA'

    # Set parameter lists containing plotting parameters
    # param = [autoPlot, plotIntensity, plotAverage, plotNumberOfExp, plotFileName, 
    #          plotTitle, plotLegend, colors, cmap, font, showTitle]
    param = [0, 0, 1, 0, plotName, plotLegend, [0], cmap, font]

    [autoPlot, plotIntData, avgPlot, plotNumberOfExp, plotFileName, plotLegend, colors, cmap, font] = param

    # Preliminary parameters
    corrLimX = [0.5, 1e5]
    corrLimY = [-0.05, 1.05]
    intLimX = [0.4e-3, 10]
    intLimY = [-2,27]

    corrTextCoord = [2*corrLimX[0], 0.92*corrLimY[1]]
    intTextCoord = [2*intLimX[0], 0.92*intLimY[1]] 

    text = 'Number of Experiments: {:.0f}'

    multipleFiles = False

    # First time parameter (true only for first loop iteration)
    firstTime = True

    # Array to save whether intensity data is stored in each given file
    intMeasured = np.zeros(len(fileNames))

    # Preliminary check for intensity data
    for file, i in zip(fileNames, range(len(fileNames))):

        if intData[file].any():
            # Set bool to true
            intMeasured[i] = 1

    if len(fileNames) > 1:
        multipleFiles = True


    # Loop over all given files
    for i, file in enumerate(fileNames):

        # Check shape of correlation data array
        corrShape = np.shape(corrData[file])

        # Calculate number of experiments (one triplicate per experiment)
        corrExpNum = int((corrShape[1] - 1)/3)

        # Adjust colormap settings
        if multipleFiles:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,len(fileNames))))
        else:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,corrExpNum)))

        # Extract correlation data triplets (and compute row-wise mean & std. dev.)
        lagTime = corrData[file][:, 0]

        meanExpCorr = np.empty((corrShape[0], int(corrExpNum)))
        stdExpCorr = np.empty((corrShape[0], int(corrExpNum)))

        for j in range(int(corrExpNum)): 
            meanExpCorr[:,j] = np.mean(corrData[file][:, (3*j+1):(3*j+4)], axis=1)
            stdExpCorr[:,j] = np.std(corrData[file][:, (3*j+1):(3*j+4)], axis=1, ddof=1)

        if avgPlot and corrExpNum > 1:
            meanCorr = np.mean(meanExpCorr, axis=1)
            stdCorr = np.std(meanExpCorr, axis=1)

        elif avgPlot and corrExpNum == 1:
            meanCorr = meanExpCorr[:,0]
            stdCorr = stdExpCorr[:,0]

        else:
            meanCorr = meanExpCorr
            stdCorr = stdExpCorr


        if firstTime:
            
            # Skip first color
            ax3.plot([],[], label='_nolegend_')
            ax3.fill_between([], [], [], label='_nolegend_')

            # Set axis labels and legends
            ax3.set_xlabel('Lag Time [$\mu$s]')
            #ax3.set_ylabel('Correlation Coefficient [-]')
            ax3.set_xlim(corrLimX)
            ax3.set_ylim(corrLimY)
            ax3.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                                bottom=True, top=False, left=True, right=True, direction='in')
            if multipleFiles and plotNumberOfExp:
                ax3.text(corrTextCoord[0], corrTextCoord[1], text.format(len(fileNames)), va='center')
            elif plotNumberOfExp:
                ax3.text(corrTextCoord[0], corrTextCoord[1], text.format(corrExpNum), va='center')

            firstTime = False

        if avgPlot:
            # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
            if colors[0]:
                ax3.semilogx(lagTime, meanCorr, color=colors[i])
                ax3.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, color=colors[i], alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)
            else:
                ax3.semilogx(lagTime, meanCorr)
                ax3.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)

            ax3.xaxis.get_major_locator().set_params(numticks=99)
            ax3.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')


        else:
            for j in range(int(corrExpNum)): 
                # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
                if colors[0]:
                    ax3.semilogx(lagTime, meanCorr[:,j], color=colors[i])
                    ax3.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                    color=colors[i], alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                    antialiased=True)
                else:
                    if multipleFiles:
                        fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                        ax3.semilogx(lagTime, meanCorr[:,j], color=fileColor)
                        ax3.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        color=fileColor, alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)
                    else:
                        ax3.semilogx(lagTime, meanCorr[:,j])
                        ax3.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)

                ax3.xaxis.get_major_locator().set_params(numticks=99)
                ax3.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')

        # Skip one color
        ax3.plot([],[], label='_nolegend_')
        ax3.fill_between([], [], [], label='_nolegend_')

        #if plotLegend and i == (len(fileNames) - 1):
        #    ax3.legend(plotLegend)

        if not avgPlot:
            for i in range(len(fileNames)):
                leg = ax3.get_legend()
                fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                leg.legendHandles[i].set_color(fileColor)

    plotLegend = ['Sheared', 'Sheared, Air', 'Unsheared']
    fileNames = ['JJ1_50mg_200rpm_10-00min', 'JJ1_50mg_AIR_010%v_200rpm_10-00min', 'JJ1_50mg_000rpm_00-00min_NA']
    
    fileNames = [fileName + '_' + mat for fileName in fileNames]

    fileNames[2] = 'JJ1_50mg_000rpm_00-00min_NA'

    # Set parameter lists containing plotting parameters
    # param = [autoPlot, plotIntensity, plotAverage, plotNumberOfExp, plotFileName, 
    #          plotTitle, plotLegend, colors, cmap, font, showTitle]
    param = [0, 0, 1, 0, plotName, plotLegend, [0], cmap, font]

    [autoPlot, plotIntData, avgPlot, plotNumberOfExp, plotFileName, plotLegend, colors, cmap, font] = param

    # Preliminary parameters
    corrLimX = [0.5, 1e5]
    corrLimY = [-0.05, 1.05]
    intLimX = [0.4e-3, 10]
    intLimY = [-2,27]

    corrTextCoord = [2*corrLimX[0], 0.92*corrLimY[1]]
    intTextCoord = [2*intLimX[0], 0.92*intLimY[1]] 

    text = 'Number of Experiments: {:.0f}'

    multipleFiles = False

    # First time parameter (true only for first loop iteration)
    firstTime = True

    # Array to save whether intensity data is stored in each given file
    intMeasured = np.zeros(len(fileNames))

    # Preliminary check for intensity data
    for file, i in zip(fileNames, range(len(fileNames))):

        if intData[file].any():
            # Set bool to true
            intMeasured[i] = 1

    if len(fileNames) > 1:
        multipleFiles = True


    # Loop over all given files
    for i, file in enumerate(fileNames):

        # Check shape of correlation data array
        corrShape = np.shape(corrData[file])

        # Calculate number of experiments (one triplicate per experiment)
        corrExpNum = int((corrShape[1] - 1)/3)

        # Adjust colormap settings
        if multipleFiles:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,len(fileNames))))
        else:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,corrExpNum)))

        # Extract correlation data triplets (and compute row-wise mean & std. dev.)
        lagTime = corrData[file][:, 0]

        meanExpCorr = np.empty((corrShape[0], int(corrExpNum)))
        stdExpCorr = np.empty((corrShape[0], int(corrExpNum)))

        for j in range(int(corrExpNum)): 
            meanExpCorr[:,j] = np.mean(corrData[file][:, (3*j+1):(3*j+4)], axis=1)
            stdExpCorr[:,j] = np.std(corrData[file][:, (3*j+1):(3*j+4)], axis=1, ddof=1)

        if avgPlot and corrExpNum > 1:
            meanCorr = np.mean(meanExpCorr, axis=1)
            stdCorr = np.std(meanExpCorr, axis=1)

        elif avgPlot and corrExpNum == 1:
            meanCorr = meanExpCorr[:,0]
            stdCorr = stdExpCorr[:,0]

        else:
            meanCorr = meanExpCorr
            stdCorr = stdExpCorr


        if firstTime:
            
            # Skip first color
            ax4.plot([],[], label='_nolegend_')
            ax4.fill_between([], [], [], label='_nolegend_')

            # Set axis labels and legends
            ax4.set_xlabel('Lag Time [$\mu$s]')
            #ax4.set_ylabel('Correlation Coefficient [-]')
            ax4.set_xlim(corrLimX)
            ax4.set_ylim(corrLimY)
            ax4.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                                bottom=True, top=False, left=True, right=True, direction='in')
            if multipleFiles and plotNumberOfExp:
                ax4.text(corrTextCoord[0], corrTextCoord[1], text.format(len(fileNames)), va='center')
            elif plotNumberOfExp:
                ax4.text(corrTextCoord[0], corrTextCoord[1], text.format(corrExpNum), va='center')

            firstTime = False

        if avgPlot:
            # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
            if colors[0]:
                ax4.semilogx(lagTime, meanCorr, color=colors[i])
                ax4.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, color=colors[i], alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)
            else:
                ax4.semilogx(lagTime, meanCorr)
                ax4.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)

            ax4.xaxis.get_major_locator().set_params(numticks=99)
            ax4.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')


        else:
            for j in range(int(corrExpNum)): 
                # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
                if colors[0]:
                    ax4.semilogx(lagTime, meanCorr[:,j], color=colors[i])
                    ax4.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                    color=colors[i], alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                    antialiased=True)
                else:
                    if multipleFiles:
                        fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                        ax4.semilogx(lagTime, meanCorr[:,j], color=fileColor)
                        ax4.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        color=fileColor, alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)
                    else:
                        ax4.semilogx(lagTime, meanCorr[:,j])
                        ax4.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)

                ax4.xaxis.get_major_locator().set_params(numticks=99)
                ax4.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')

        # Skip one color
        ax4.plot([],[], label='_nolegend_')
        ax4.fill_between([], [], [], label='_nolegend_')

        #if plotLegend and i == (len(fileNames) - 1):
        #    ax4.legend(plotLegend, bbox_to_anchor=(1, 1.05))

        if not avgPlot:
            for i in range(len(fileNames)):
                leg = ax4.get_legend()
                fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                leg.legendHandles[i].set_color(fileColor)
                    
    #endNamePdf = plotFileName + '.pdf'
    #endNamePng = plotFileName + '.png'

    #if logPlot:
    #    exportNamePdf = os.path.join('MDI', 'Overview', 'Log', 'PDF', endNamePdf)
    #    exportNamePng = os.path.join('MDI', 'Overview', 'Log', 'PNG',endNamePng)
    #else:
    #    exportNamePdf = os.path.join('MDI', 'Overview', 'Linear', 'PDF', endNamePdf)
    #    exportNamePng = os.path.join('MDI', 'Overview', 'Linear', 'PNG',endNamePng)

    #if showTitle:                                 
        # Set overall title
        #fig.suptitle(plotTitle, y=1.06)

    # Save plot    
    if mat == 'PM':   
        plt.savefig(os.path.join('02_Plots', '00_Summary', '01_PDF', 'Air_PM.pdf'), bbox_inches='tight')
        plt.savefig(os.path.join('02_Plots', '00_Summary', '02_PNG', 'Air_PM.png'), bbox_inches='tight')
    else:
        plt.savefig(os.path.join('02_Plots', '00_Summary', '01_PDF', 'Air_ST.pdf'), bbox_inches='tight')
        plt.savefig(os.path.join('02_Plots', '00_Summary', '02_PNG', 'Air_ST.png'), bbox_inches='tight')

    # Don't show plot in jupyter notebook
    #plt.close()

def plotPS2ConcPub(corrData, intData, countData, plotName, cmap, font, mat):

    """
    Plot MDI and DLS results in a subplot matrix
    """

    ######## MDI Plot ########    

    expNames = ('Sheared', 'Sheared, PS 20', 'Unsheared', 
                'Sheared', 'Sheared, PS 20', 'Unsheared', 
                'Sheared', 'Sheared, PS 20', 'Unsheared')

    fileNames = ['JJ1_01mg_200rpm_10-00min', 'JJ1_01mg_PS2_250ug_200rpm_10-00min', 'JJ1_01mg_000rpm_00-00min_NA',
                 'JJ1_10mg_200rpm_10-00min', 'JJ1_10mg_PS2_250ug_200rpm_10-00min', 'JJ1_10mg_000rpm_00-00min_NA',
                 'JJ1_50mg_200rpm_10-00min', 'JJ1_50mg_PS2_500ug_200rpm_10-00min', 'JJ1_50mg_000rpm_00-00min_NA']
    
    fileNames = [fileName + '_' + mat for fileName in fileNames]
    fileNames[2] = 'JJ1_01mg_000rpm_00-00min_NA'
    fileNames[5] = 'JJ1_10mg_000rpm_00-00min_NA'
    fileNames[8] = 'JJ1_50mg_000rpm_00-00min_NA'


    # Set parameter lists containing plotting parameters
    # param = [autoPlot, rmFirstBin, logPlot, plotAverage, plotLimits, expNames, plotFileName, 
    #          plotTitle, colors, cmap, font, showTitle]
    param = [0, 0, 0, 1, (10, 4e6), expNames, plotName, [0], cmap, font]

    [autoPlot, rmFirstBin, logPlot, avgPlot, plotLimits, expNames, plotFileName, colors, cmap, font] = param

    # Preliminary parameters
    countLimY = plotLimits

    # Tuple containing experiment number strings
    expTuple = expNames

    numFiles = len(fileNames)

    # Initialize empty count dictionary (needed for grouped barplot)
    keyList = ['Total ($\geq$ 1 $\mu$m)', '1-2 $\mu$m', '2-5 $\mu$m', '5-10 $\mu$m']#, '10-25 $\mu$m']
    keyLength = len(keyList)
    keyContent = []
    meanCountDict = dict.fromkeys(keyList, keyContent)
    stdCountDict = dict.fromkeys(keyList, keyContent)

    # Adjust colormap settings
    #plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,6)))
    darkGray = [0.4, 0.4, 0.4, 1]   # RGBA value, in closed interval [0,1]
    colMap = getattr(plt.cm, cmap)(np.linspace(0,1,6))
    plt.rcParams["axes.prop_cycle"] = plt.cycler("color", np.insert(colMap, 0, darkGray, axis=0))

    countDataProc = {}
    countDataDel = {}

    firstTime = True

    # Loop over all given files
    for i, file in enumerate(fileNames):

        countDataProc[file] = np.insert(countData[file], 0, np.sum(countData[file], axis=0), axis=0)

        # Find max. count in given file
        countMax = np.max(countDataProc[file])

        # Check shape of correlation data array
        countShape = np.shape(countDataProc[file])

        # Calculate number of experiments (one duplicate per experiment)
        countExpNum = int(countShape[1] / 2)

        if firstTime:
            meanExpCount = np.zeros((countShape[0], countExpNum*numFiles))
            stdExpCount = np.zeros((countShape[0], countExpNum*numFiles)) 

        for j in range(int(countExpNum)):
            meanExpCount[:,i*countExpNum+j] = np.mean(countDataProc[file][:, (2*j):(2*j+2)], axis=1)
            stdExpCount[:,i*countExpNum+j] = np.std(countDataProc[file][:, (2*j):(2*j+2)], axis=1)#, ddof=1)

        if avgPlot:

            if firstTime:
                meanCount = np.zeros((countShape[0], numFiles))
                stdCount = np.zeros((countShape[0], numFiles))

            if i == 0:
                meanCount[:,i] = np.mean(meanExpCount[:,:(i*countExpNum+countExpNum)], axis=1)
                stdCount[:,i] = np.std(meanExpCount[:,:(i*countExpNum+countExpNum)], axis=1)#, ddof=1)
            else: 
                meanCount[:,i] = np.mean(meanExpCount[:,((i-1)*countExpNum+countExpNum):(i*countExpNum+countExpNum)], axis=1)
                stdCount[:,i] = np.std(meanExpCount[:,((i-1)*countExpNum+countExpNum):(i*countExpNum+countExpNum)], axis=1)#, ddof=1)

        if firstTime:

            figWidth = 20
            figHeight = figWidth/1.8
            
            fig = plt.figure(figsize=(figWidth, figHeight))
            gs = fig.add_gridspec(2,5, hspace=0.3, width_ratios=(0.5, 3, 3, 3, 0.1), height_ratios=(1.5, 1))
            ax = fig.add_subplot(gs[0, :])
            ax2 = fig.add_subplot(gs[1, 1])
            ax3 = fig.add_subplot(gs[1, 2])
            ax4 = fig.add_subplot(gs[1, 3])


            # Add some text for labels, title and custom x-axis tick labels, etc.
            #ax.set_xlabel('Experiments')
            ax.set_ylabel('Number Density [mL$^{-1}$]')
            ax.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                            bottom=True, top=False, left=True, right=True, direction='in')

            firstTime = False


    x = np.arange(len(expTuple))
    width = 0.2  # the width of the bars
    multiplier = 1

    for j in range(keyLength):                
        meanCountDict[keyList[j]] = meanCount[j,:]
        stdCountDict[keyList[j]] = stdCount[j,:]
        

    for attribute, measurement in meanCountDict.items():
        offset = width * multiplier
        rects = ax.bar(x + offset, measurement, width, label=attribute, alpha=0.7)
        ax.errorbar(x + offset, measurement, yerr=stdCountDict[attribute], ecolor='black',
                capsize=5, linestyle='')
        multiplier += 1
            

    # Set additional plotting parameters
    ax.set_xticks(x + (keyLength+1)/2*width, expTuple)
    ax.legend(ncols=1)

    if logPlot:
        ax.set_yscale('log')
    else:
        ax.set_yscale('linear')

    if plotLimits and logPlot:
        ax.set_ylim(countLimY)
    elif plotLimits:
        ax.set_ylim(countLimY)
        ax.ticklabel_format(style='sci', axis='y', scilimits=(0,0), useMathText=True)
    elif logPlot:
        ax.set_ylim(bottom=10)
    else:
        ax.set_ylim(bottom=0)
        ax.ticklabel_format(style='sci', axis='y', scilimits=(0,0), useMathText=True)

    ######## DLS Plots ########    

    plotLegend = ['Sheared', 'Sheared, PS 20', 'Unsheared']
    fileNames = ['JJ1_01mg_200rpm_10-00min', 'JJ1_01mg_PS2_250ug_200rpm_10-00min', 'JJ1_01mg_000rpm_00-00min_NA']

    fileNames = [fileName + '_' + mat for fileName in fileNames]

    fileNames[2] = 'JJ1_01mg_000rpm_00-00min_NA'


    # Set parameter lists containing plotting parameters
    # param = [autoPlot, plotIntensity, plotAverage, plotNumberOfExp, plotFileName, 
    #          plotTitle, plotLegend, colors, cmap, font, showTitle]
    param = [0, 0, 1, 0, plotName, plotLegend, [0], cmap, font]

    [autoPlot, plotIntData, avgPlot, plotNumberOfExp, plotFileName, plotLegend, colors, cmap, font] = param

    # Preliminary parameters
    corrLimX = [0.5, 1e5]
    corrLimY = [-0.05, 1.05]
    intLimX = [0.4e-3, 10]
    intLimY = [-2,27]

    corrTextCoord = [2*corrLimX[0], 0.92*corrLimY[1]]
    intTextCoord = [2*intLimX[0], 0.92*intLimY[1]] 

    text = 'Number of Experiments: {:.0f}'

    multipleFiles = False

    # First time parameter (true only for first loop iteration)
    firstTime = True

    # Array to save whether intensity data is stored in each given file
    intMeasured = np.zeros(len(fileNames))

    # Preliminary check for intensity data
    for file, i in zip(fileNames, range(len(fileNames))):

        if intData[file].any():
            # Set bool to true
            intMeasured[i] = 1

    if len(fileNames) > 1:
        multipleFiles = True

    # Loop over all given files
    for i, file in enumerate(fileNames):

        # Check shape of correlation data array
        corrShape = np.shape(corrData[file])

        # Calculate number of experiments (one triplicate per experiment)
        corrExpNum = int((corrShape[1] - 1)/3)

        # Adjust colormap settings
        if multipleFiles:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,len(fileNames))))
        else:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,corrExpNum)))

        # Extract correlation data triplets (and compute row-wise mean & std. dev.)
        lagTime = corrData[file][:, 0]

        meanExpCorr = np.empty((corrShape[0], int(corrExpNum)))
        stdExpCorr = np.empty((corrShape[0], int(corrExpNum)))

        for j in range(int(corrExpNum)): 
            meanExpCorr[:,j] = np.mean(corrData[file][:, (3*j+1):(3*j+4)], axis=1)
            stdExpCorr[:,j] = np.std(corrData[file][:, (3*j+1):(3*j+4)], axis=1, ddof=1)

        if avgPlot and corrExpNum > 1:
            meanCorr = np.mean(meanExpCorr, axis=1)
            stdCorr = np.std(meanExpCorr, axis=1)

        elif avgPlot and corrExpNum == 1:
            meanCorr = meanExpCorr[:,0]
            stdCorr = stdExpCorr[:,0]

        else:
            meanCorr = meanExpCorr
            stdCorr = stdExpCorr


        if firstTime:
            # Skip first color
            ax2.plot([],[], label='_nolegend_')
            ax2.fill_between([], [], [], label='_nolegend_')
            # Set axis labels and legends
            ax2.set_xlabel('Lag Time [$\mu$s]')
            ax2.set_ylabel('Correlation Coefficient [-]')
            ax2.set_xlim(corrLimX)
            ax2.set_ylim(corrLimY)
            ax2.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                                bottom=True, top=False, left=True, right=True, direction='in')
            if multipleFiles and plotNumberOfExp:
                ax2.text(corrTextCoord[0], corrTextCoord[1], text.format(len(fileNames)), va='center')
            elif plotNumberOfExp:
                ax2.text(corrTextCoord[0], corrTextCoord[1], text.format(corrExpNum), va='center')
                
            firstTime = False
            

        if avgPlot:
            # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
            if colors[0]:
                ax2.semilogx(lagTime, meanCorr, color=colors[i])
                ax2.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, color=colors[i], alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)
            else:
                ax2.semilogx(lagTime, meanCorr)
                ax2.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)

            ax2.xaxis.get_major_locator().set_params(numticks=99)
            ax2.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')


        else:
            for j in range(int(corrExpNum)): 
                # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
                if colors[0]:
                    ax2.semilogx(lagTime, meanCorr[:,j], color=colors[i])
                    ax2.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                    color=colors[i], alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                    antialiased=True)
                else:
                    if multipleFiles:
                        fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                        ax2.semilogx(lagTime, meanCorr[:,j], color=fileColor)
                        ax2.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        color=fileColor, alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)
                    else:
                        ax2.semilogx(lagTime, meanCorr[:,j])
                        ax2.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)

                ax2.xaxis.get_major_locator().set_params(numticks=99)
                ax2.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')

        # Skip one color
        ax2.plot([],[], label='_nolegend_')
        ax2.fill_between([], [], [], label='_nolegend_')

        if plotLegend and i == (len(fileNames) - 1):
            ax2.legend(plotLegend, bbox_to_anchor=(-0.25, 1.05))

            if not avgPlot:
                for i in range(len(fileNames)):
                    leg = ax2.get_legend()
                    fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                    leg.legendHandles[i].set_color(fileColor)

    plotLegend = ['Sheared', 'Sheared, PS 20', 'Unsheared']
    fileNames = ['JJ1_10mg_200rpm_10-00min', 'JJ1_10mg_PS2_250ug_200rpm_10-00min', 'JJ1_10mg_000rpm_00-00min_NA']
    
    fileNames = [fileName + '_' + mat for fileName in fileNames]

    fileNames[2] = 'JJ1_10mg_000rpm_00-00min_NA'

    # Set parameter lists containing plotting parameters
    # param = [autoPlot, plotIntensity, plotAverage, plotNumberOfExp, plotFileName, 
    #          plotTitle, plotLegend, colors, cmap, font, showTitle]
    param = [0, 0, 1, 0, plotName, plotLegend, [0], cmap, font]

    [autoPlot, plotIntData, avgPlot, plotNumberOfExp, plotFileName, plotLegend, colors, cmap, font] = param

    # Preliminary parameters
    corrLimX = [0.5, 1e5]
    corrLimY = [-0.05, 1.05]
    intLimX = [0.4e-3, 10]
    intLimY = [-2,27]

    corrTextCoord = [2*corrLimX[0], 0.92*corrLimY[1]]
    intTextCoord = [2*intLimX[0], 0.92*intLimY[1]] 

    text = 'Number of Experiments: {:.0f}'

    multipleFiles = False

    # First time parameter (true only for first loop iteration)
    firstTime = True

    # Array to save whether intensity data is stored in each given file
    intMeasured = np.zeros(len(fileNames))

    # Preliminary check for intensity data
    for file, i in zip(fileNames, range(len(fileNames))):

        if intData[file].any():
            # Set bool to true
            intMeasured[i] = 1

    if len(fileNames) > 1:
        multipleFiles = True


    # Loop over all given files
    for i, file in enumerate(fileNames):

        # Check shape of correlation data array
        corrShape = np.shape(corrData[file])

        # Calculate number of experiments (one triplicate per experiment)
        corrExpNum = int((corrShape[1] - 1)/3)

        # Adjust colormap settings
        if multipleFiles:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,len(fileNames))))
        else:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,corrExpNum)))

        # Extract correlation data triplets (and compute row-wise mean & std. dev.)
        lagTime = corrData[file][:, 0]

        meanExpCorr = np.empty((corrShape[0], int(corrExpNum)))
        stdExpCorr = np.empty((corrShape[0], int(corrExpNum)))

        for j in range(int(corrExpNum)): 
            meanExpCorr[:,j] = np.mean(corrData[file][:, (3*j+1):(3*j+4)], axis=1)
            stdExpCorr[:,j] = np.std(corrData[file][:, (3*j+1):(3*j+4)], axis=1, ddof=1)

        if avgPlot and corrExpNum > 1:
            meanCorr = np.mean(meanExpCorr, axis=1)
            stdCorr = np.std(meanExpCorr, axis=1)

        elif avgPlot and corrExpNum == 1:
            meanCorr = meanExpCorr[:,0]
            stdCorr = stdExpCorr[:,0]

        else:
            meanCorr = meanExpCorr
            stdCorr = stdExpCorr


        if firstTime:
            
            # Skip first color
            ax3.plot([],[], label='_nolegend_')
            ax3.fill_between([], [], [], label='_nolegend_')

            # Set axis labels and legends
            ax3.set_xlabel('Lag Time [$\mu$s]')
            #ax3.set_ylabel('Correlation Coefficient [-]')
            ax3.set_xlim(corrLimX)
            ax3.set_ylim(corrLimY)
            ax3.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                                bottom=True, top=False, left=True, right=True, direction='in')
            if multipleFiles and plotNumberOfExp:
                ax3.text(corrTextCoord[0], corrTextCoord[1], text.format(len(fileNames)), va='center')
            elif plotNumberOfExp:
                ax3.text(corrTextCoord[0], corrTextCoord[1], text.format(corrExpNum), va='center')

            firstTime = False

        if avgPlot:
            # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
            if colors[0]:
                ax3.semilogx(lagTime, meanCorr, color=colors[i])
                ax3.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, color=colors[i], alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)
            else:
                ax3.semilogx(lagTime, meanCorr)
                ax3.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)

            ax3.xaxis.get_major_locator().set_params(numticks=99)
            ax3.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')


        else:
            for j in range(int(corrExpNum)): 
                # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
                if colors[0]:
                    ax3.semilogx(lagTime, meanCorr[:,j], color=colors[i])
                    ax3.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                    color=colors[i], alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                    antialiased=True)
                else:
                    if multipleFiles:
                        fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                        ax3.semilogx(lagTime, meanCorr[:,j], color=fileColor)
                        ax3.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        color=fileColor, alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)
                    else:
                        ax3.semilogx(lagTime, meanCorr[:,j])
                        ax3.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)

                ax3.xaxis.get_major_locator().set_params(numticks=99)
                ax3.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')

        # Skip one color
        ax3.plot([],[], label='_nolegend_')
        ax3.fill_between([], [], [], label='_nolegend_')

        #if plotLegend and i == (len(fileNames) - 1):
        #    ax3.legend(plotLegend)

        if not avgPlot:
            for i in range(len(fileNames)):
                leg = ax3.get_legend()
                fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                leg.legendHandles[i].set_color(fileColor)

    plotLegend = ['Sheared', 'Sheared, PS 20', 'Unsheared']
    fileNames = ['JJ1_50mg_200rpm_10-00min', 'JJ1_50mg_PS2_500ug_200rpm_10-00min', 'JJ1_50mg_000rpm_00-00min_NA']
    
    fileNames = [fileName + '_' + mat for fileName in fileNames]

    fileNames[2] = 'JJ1_50mg_000rpm_00-00min_NA'

    # Set parameter lists containing plotting parameters
    # param = [autoPlot, plotIntensity, plotAverage, plotNumberOfExp, plotFileName, 
    #          plotTitle, plotLegend, colors, cmap, font, showTitle]
    param = [0, 0, 1, 0, plotName, plotLegend, [0], cmap, font]

    [autoPlot, plotIntData, avgPlot, plotNumberOfExp, plotFileName, plotLegend, colors, cmap, font] = param

    # Preliminary parameters
    corrLimX = [0.5, 1e5]
    corrLimY = [-0.05, 1.05]
    intLimX = [0.4e-3, 10]
    intLimY = [-2,27]

    corrTextCoord = [2*corrLimX[0], 0.92*corrLimY[1]]
    intTextCoord = [2*intLimX[0], 0.92*intLimY[1]] 

    text = 'Number of Experiments: {:.0f}'

    multipleFiles = False

    # First time parameter (true only for first loop iteration)
    firstTime = True

    # Array to save whether intensity data is stored in each given file
    intMeasured = np.zeros(len(fileNames))

    # Preliminary check for intensity data
    for file, i in zip(fileNames, range(len(fileNames))):

        if intData[file].any():
            # Set bool to true
            intMeasured[i] = 1

    if len(fileNames) > 1:
        multipleFiles = True


    # Loop over all given files
    for i, file in enumerate(fileNames):

        # Check shape of correlation data array
        corrShape = np.shape(corrData[file])

        # Calculate number of experiments (one triplicate per experiment)
        corrExpNum = int((corrShape[1] - 1)/3)

        # Adjust colormap settings
        if multipleFiles:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,len(fileNames))))
        else:
            plt.rcParams["axes.prop_cycle"] = plt.cycler("color", getattr(plt.cm, cmap)(np.linspace(0,1,corrExpNum)))

        # Extract correlation data triplets (and compute row-wise mean & std. dev.)
        lagTime = corrData[file][:, 0]

        meanExpCorr = np.empty((corrShape[0], int(corrExpNum)))
        stdExpCorr = np.empty((corrShape[0], int(corrExpNum)))

        for j in range(int(corrExpNum)): 
            meanExpCorr[:,j] = np.mean(corrData[file][:, (3*j+1):(3*j+4)], axis=1)
            stdExpCorr[:,j] = np.std(corrData[file][:, (3*j+1):(3*j+4)], axis=1, ddof=1)

        if avgPlot and corrExpNum > 1:
            meanCorr = np.mean(meanExpCorr, axis=1)
            stdCorr = np.std(meanExpCorr, axis=1)

        elif avgPlot and corrExpNum == 1:
            meanCorr = meanExpCorr[:,0]
            stdCorr = stdExpCorr[:,0]

        else:
            meanCorr = meanExpCorr
            stdCorr = stdExpCorr


        if firstTime:
            
            # Skip first color
            ax4.plot([],[], label='_nolegend_')
            ax4.fill_between([], [], [], label='_nolegend_')

            # Set axis labels and legends
            ax4.set_xlabel('Lag Time [$\mu$s]')
            #ax4.set_ylabel('Correlation Coefficient [-]')
            ax4.set_xlim(corrLimX)
            ax4.set_ylim(corrLimY)
            ax4.tick_params(labelbottom=True, labeltop=False, labelleft=True, labelright=False,
                                bottom=True, top=False, left=True, right=True, direction='in')
            if multipleFiles and plotNumberOfExp:
                ax4.text(corrTextCoord[0], corrTextCoord[1], text.format(len(fileNames)), va='center')
            elif plotNumberOfExp:
                ax4.text(corrTextCoord[0], corrTextCoord[1], text.format(corrExpNum), va='center')

            firstTime = False

        if avgPlot:
            # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
            if colors[0]:
                ax4.semilogx(lagTime, meanCorr, color=colors[i])
                ax4.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, color=colors[i], alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)
            else:
                ax4.semilogx(lagTime, meanCorr)
                ax4.fill_between(lagTime, meanCorr-stdCorr, meanCorr+stdCorr, alpha=0.25,
                linewidth=1, linestyle='solid', label='_nolegend_', antialiased=True)

            ax4.xaxis.get_major_locator().set_params(numticks=99)
            ax4.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')


        else:
            for j in range(int(corrExpNum)): 
                # Add curves to subfigure 1 (plot correlation [-] vs. lagTime [us])
                if colors[0]:
                    ax4.semilogx(lagTime, meanCorr[:,j], color=colors[i])
                    ax4.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                    color=colors[i], alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                    antialiased=True)
                else:
                    if multipleFiles:
                        fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                        ax4.semilogx(lagTime, meanCorr[:,j], color=fileColor)
                        ax4.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        color=fileColor, alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)
                    else:
                        ax4.semilogx(lagTime, meanCorr[:,j])
                        ax4.fill_between(lagTime, meanCorr[:,j]-stdCorr[:,j], meanCorr[:,j]+stdCorr[:,j], 
                        alpha=0.25, linewidth=1, linestyle='solid', label='_nolegend_', 
                        antialiased=True)

                ax4.xaxis.get_major_locator().set_params(numticks=99)
                ax4.xaxis.get_minor_locator().set_params(numticks=99, subs='auto')

        # Skip one color
        ax4.plot([],[], label='_nolegend_')
        ax4.fill_between([], [], [], label='_nolegend_')

        #if plotLegend and i == (len(fileNames) - 1):
        #    ax4.legend(plotLegend, bbox_to_anchor=(1, 1.05))

        if not avgPlot:
            for i in range(len(fileNames)):
                leg = ax4.get_legend()
                fileColor = plt.rcParams['axes.prop_cycle'].by_key()['color'][i]
                leg.legendHandles[i].set_color(fileColor)
                    
    #endNamePdf = plotFileName + '.pdf'
    #endNamePng = plotFileName + '.png'

    #if logPlot:
    #    exportNamePdf = os.path.join('MDI', 'Overview', 'Log', 'PDF', endNamePdf)
    #    exportNamePng = os.path.join('MDI', 'Overview', 'Log', 'PNG',endNamePng)
    #else:
    #    exportNamePdf = os.path.join('MDI', 'Overview', 'Linear', 'PDF', endNamePdf)
    #    exportNamePng = os.path.join('MDI', 'Overview', 'Linear', 'PNG',endNamePng)

    #if showTitle:                                 
        # Set overall title
        #fig.suptitle(plotTitle, y=1.06)

    # Save plot    
    if mat == 'PM':   
        plt.savefig(os.path.join('02_Plots', '00_Summary', '01_PDF', 'PS2Conc_PM.pdf'), bbox_inches='tight')
        plt.savefig(os.path.join('02_Plots', '00_Summary', '02_PNG', 'PS2Conc_PM.png'), bbox_inches='tight')
    else:
        plt.savefig(os.path.join('02_Plots', '00_Summary', '01_PDF', 'PS2Conc_ST.pdf'), bbox_inches='tight')
        plt.savefig(os.path.join('02_Plots', '00_Summary', '02_PNG', 'PS2Conc_ST.png'), bbox_inches='tight')

    # Don't show plot in jupyter notebook
    #plt.close()





def addSlide(prs, img, img2, layout, titleText, subtitleText, leftMargin, topMargin, imWidth, imHeight, isIndiv):

    slide = prs.slides.add_slide(layout)

    title = slide.shapes.title
    title.text = titleText

    txBox = slide.shapes.add_textbox(Inches(0.7), 0.6*Inches(2), width=Inches(5.5), height=Inches(5.5))
    tf = txBox.text_frame
    tf.text = subtitleText

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame    
    p = tf.add_paragraph()
    p.text = ''
    p.level = 1

    leftCentered = (prs.slide_width - imWidth) / 2

    if img2 and not isIndiv:
        slide.shapes.add_picture(img, 0.1*leftCentered, topMargin, width=imWidth, height=None)
        slide.shapes.add_picture(img2, 1.9*leftCentered, topMargin, width=imWidth, height=None)
    elif img2 and isIndiv:
        slide.shapes.add_picture(img, 0.1*leftCentered, topMargin, width=None, height=imHeight)
        slide.shapes.add_picture(img2, 1.9*leftCentered, topMargin, width=None, height=imHeight) 
    elif isIndiv:
        slide.shapes.add_picture(img, leftMargin, topMargin, width=None, height=imHeight)
    else:
        slide.shapes.add_picture(img, leftCentered, topMargin, width=imWidth, height=None)

    return slide

